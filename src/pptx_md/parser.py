"""PPTX parser: Presentation -> PresentationIR (FR-03, FR-04).

Uses only python-pptx (and stdlib). MUST NOT import VLM SDKs
(NFR-08, ADR-002). Boundary module: narrows python-pptx ``Any`` into typed IR
(§5.1, ADR-206).

Public API:
    parse_presentation(path) -> PresentationIR
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation as _PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError

from pptx_md.errors import ParseError
from pptx_md.ir import (
    GroupShapeIR,
    ImageShapeIR,
    OtherShapeIR,
    ParagraphIR,
    PresentationIR,
    ShapeIR,
    ShapeKind,
    SlideIR,
    TableShapeIR,
    TextShapeIR,
)

_log = logging.getLogger("pptx_md.parser")

# FR-26: Chart/SmartArt GraphicFrame best-effort text/type extraction.
_A_T_QN = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
_A_GRAPHICDATA_QN = "{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData"
_C_V_QN = "{http://schemas.openxmlformats.org/drawingml/2006/chart}v"
_C_STRCACHE_QN = "{http://schemas.openxmlformats.org/drawingml/2006/chart}strCache"
_DIAGRAM_GRAPHIC_DATA_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"


def parse_presentation(path: str | Path) -> PresentationIR:
    """Parse a PPTX file into a PresentationIR tree (FR-03, FR-04).

    Args:
        path: Path to the PPTX file (str or pathlib.Path).

    Returns:
        PresentationIR containing all slides and shapes in original order.

    Raises:
        ParseError: If the file does not exist or is not a valid PPTX package.
        FileNotFoundError: Passed through if raised before pptx open attempt.
    """
    p = Path(path)
    if not p.is_file():
        raise ParseError(f"file not found: {p}")

    try:
        prs = _PptxPresentation(str(p))
    except PackageNotFoundError as exc:
        raise ParseError(f"not a valid PPTX package: {p}") from exc
    except Exception as exc:
        raise ParseError(f"failed to open PPTX: {p}") from exc

    pres = PresentationIR(
        source_path=str(p),
        slide_width_emu=int(prs.slide_width or 0),
        slide_height_emu=int(prs.slide_height or 0),
    )
    for idx, slide in enumerate(prs.slides):
        pres.slides.append(_parse_slide(idx, slide))

    return pres


def _parse_slide(idx: int, slide: object) -> SlideIR:
    """Parse a single slide into SlideIR (FR-03)."""
    title = _safe_title(slide)
    notes = _safe_notes(slide)
    sld = SlideIR(index=idx, title=title, notes=notes)

    for shape in getattr(slide, "shapes", []):
        sld.shapes.append(_parse_shape(shape))

    return sld


def _extract_coords(shape: object) -> tuple[int, int, int, int]:
    """Read left/top/width/height from a shape, defaulting None to 0 (FR-23 AC3).

    python-pptx returns EMU integers for these attributes, but some shapes
    (e.g. placeholders on certain layouts) may return None.  We normalise
    None -> 0 to keep the IR schema non-nullable (ADR-202).

    Returns:
        (left, top, width, height) all as int EMU values.
    """

    def _emu(attr: str) -> int:
        val = getattr(shape, attr, None)
        if val is None:
            return 0
        try:
            return int(val)
        except (TypeError, ValueError):
            return 0

    return _emu("left"), _emu("top"), _emu("width"), _emu("height")


def _parse_shape(shape: object) -> ShapeIR:
    """Parse one shape into the appropriate ShapeIR subtype (FR-04).

    Per-shape failures are caught and demoted to OtherShapeIR (ADR-204).
    Coordinate fields (FR-23) are always populated; None -> 0 (AC3).
    """
    sid = int(getattr(shape, "shape_id", 0) or 0)
    name = str(getattr(shape, "name", "") or "")
    left, top, width, height = _extract_coords(shape)

    try:
        ir = _dispatch_shape(shape, sid, name)
    except Exception as exc:  # ADR-204: never propagate per-shape failure
        _log.warning(
            "shape id=%d name=%r parse failed (%s); demoting to OtherShapeIR",
            sid,
            name,
            type(exc).__name__,
        )
        return OtherShapeIR(
            shape_id=sid,
            name=name,
            kind=ShapeKind.OTHER,
            left=left,
            top=top,
            width=width,
            height=height,
            mso_shape_type="UNKNOWN",
            fallback_text="",
        )

    # FR-23 AC2: stamp coordinates onto every successfully-parsed IR node
    ir.left = left
    ir.top = top
    ir.width = width
    ir.height = height
    return ir


def _dispatch_shape(shape: object, sid: int, name: str) -> ShapeIR:
    """Dispatch shape to the correct parser based on type detection.

    Priority order (ARCH-M2 §8.3):
    1. GROUP  — shape_type
    2. PICTURE — shape_type
    3. TABLE  — has_table (GraphicFrame, detected before text_frame check)
    4. TEXT   — has_text_frame
    5. OTHER  — fallback
    """
    shape_type = getattr(shape, "shape_type", None)

    if shape_type == MSO_SHAPE_TYPE.GROUP:
        return _parse_group(shape, sid, name)

    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return _parse_picture(shape, sid, name)

    if getattr(shape, "has_table", False):
        return _parse_table(shape, sid, name)

    if getattr(shape, "has_text_frame", False):
        return _parse_text(shape, sid, name)

    return _parse_other(shape, sid, name)


# ---------------------------------------------------------------------------
# Shape-type specific parsers
# ---------------------------------------------------------------------------


def _parse_text(shape: object, sid: int, name: str) -> TextShapeIR:
    """Parse a text-frame shape into TextShapeIR (FR-04 AC1, FR-21)."""
    # Determine if this is a title placeholder.
    # python-pptx raises ValueError for non-placeholder shapes, so we catch it.
    is_title = False
    is_footer = False
    try:
        placeholder_format = shape.placeholder_format  # type: ignore[attr-defined]
        if placeholder_format is not None:
            ph_type = getattr(placeholder_format, "type", None)
            if ph_type is not None:
                # PP_PLACEHOLDER.CENTER_TITLE = 3, TITLE = 15
                ph_type_str = str(ph_type)
                ph_name = str(getattr(ph_type, "name", ph_type_str)).upper()
                is_title = "TITLE" in ph_type_str or "TITLE" in ph_name
                # FR-21: FOOTER(11), SLIDE_NUMBER(12), DATE(10/13)
                is_footer = any(
                    k in ph_name for k in ("FOOTER", "SLIDE_NUMBER", "DATE")
                )
    except (ValueError, AttributeError):
        # Not a placeholder shape — is_title/is_footer stay False
        pass

    paragraphs: list[ParagraphIR] = []
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is not None:
        for para in getattr(text_frame, "paragraphs", []):
            para_text = _extract_paragraph_text(para)
            level = int(getattr(para, "level", 0) or 0)
            paragraphs.append(ParagraphIR(text=para_text, level=level))

    return TextShapeIR(
        shape_id=sid,
        name=name,
        kind=ShapeKind.TEXT,
        paragraphs=paragraphs,
        is_title=is_title,
        is_footer=is_footer,
    )


def _parse_table(shape: object, sid: int, name: str) -> TableShapeIR:
    """Parse a table shape into TableShapeIR (FR-04 AC2)."""
    table = getattr(shape, "table", None)
    rows: list[list[str]] = []

    if table is not None:
        for row in getattr(table, "rows", []):
            row_cells: list[str] = []
            for cell in getattr(row, "cells", []):
                cell_text = str(getattr(cell, "text", "") or "")
                row_cells.append(cell_text)
            rows.append(row_cells)

    n_rows = len(rows)
    n_cols = max((len(r) for r in rows), default=0)

    return TableShapeIR(
        shape_id=sid,
        name=name,
        kind=ShapeKind.TABLE,
        rows=rows,
        n_rows=n_rows,
        n_cols=n_cols,
    )


def _parse_picture(shape: object, sid: int, name: str) -> ImageShapeIR:
    """Parse a picture shape into ImageShapeIR (FR-04 AC3).

    image.blob is loaded as-is; decoding/classification is deferred to M3
    (§5.5, ADR-205).
    """
    image_bytes = b""
    image_format = ""
    image_ext = ""

    image_obj = getattr(shape, "image", None)
    if image_obj is not None:
        try:
            image_bytes = bytes(getattr(image_obj, "blob", b"") or b"")
        except Exception as exc:  # noqa: BLE001
            _log.warning("image blob read failed (%s)", type(exc).__name__)
            # Demote blob read failure to empty bytes; shape stays ImageShapeIR
            image_bytes = b""

        content_type: str = str(getattr(image_obj, "content_type", "") or "")
        # Normalise to short format label: "image/png" -> "png"
        if "/" in content_type:
            image_format = content_type.split("/")[-1].lower()
        else:
            image_format = content_type.lower()

        raw_ext: str = str(getattr(image_obj, "ext", "") or "")
        image_ext = raw_ext.lstrip(".").lower() if raw_ext else image_format

    # Alt text (§5.4, ADR-202: "" not None)
    alt_text = ""
    element = getattr(shape, "element", None)
    if element is not None:
        try:
            # python-pptx stores alt text in nvPicPr/cNvPr[@descr]
            nvPicPr = element.find(
                ".//{http://schemas.openxmlformats.org/presentationml/2006/main}nvPicPr"
            )
            if nvPicPr is not None:
                cNvPr = nvPicPr.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr"
                )
                if cNvPr is not None:
                    alt_text = str(cNvPr.get("descr", "") or "")
        except Exception:  # noqa: BLE001
            alt_text = ""

    return ImageShapeIR(
        shape_id=sid,
        name=name,
        kind=ShapeKind.IMAGE,
        image_bytes=image_bytes,
        image_format=image_format,
        image_ext=image_ext,
        alt_text=alt_text,
        classification=None,  # ADR-205: v1 always None
        description=None,  # ADR-205: v1 always None
    )


def _parse_group(shape: object, sid: int, name: str) -> GroupShapeIR:
    """Parse a group shape into GroupShapeIR, recursively (FR-04 AC4, ADR-203)."""
    children: list[ShapeIR] = []
    for child in getattr(shape, "shapes", []):
        children.append(_parse_shape(child))

    return GroupShapeIR(
        shape_id=sid,
        name=name,
        kind=ShapeKind.GROUP,
        children=children,
    )


def _parse_other(shape: object, sid: int, name: str) -> OtherShapeIR:
    """Fallback for unrecognised/unsupported shapes (FR-04 AC6, ADR-204).

    Chart/SmartArt GraphicFrame shapes are best-effort enriched (FR-26 AC5/AC6):
    ``mso_shape_type`` carries the real MSO/graphic-frame label instead of
    "UNKNOWN" when it can be determined, and ``fallback_text`` captures any
    recoverable text (chart titles, cached category/series labels, or inline
    diagram text runs) instead of always being "".
    """
    shape_type = getattr(shape, "shape_type", None)
    if shape_type is not None:
        mso_label = str(getattr(shape_type, "name", str(shape_type)))
    else:
        mso_label = _graphic_frame_type_label(shape) or "UNKNOWN"

    # Best-effort text extraction from any text_frame that might be present
    fallback_text = ""
    if getattr(shape, "has_text_frame", False):
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is not None:
            try:
                fallback_text = str(getattr(text_frame, "text", "") or "")
            except Exception:  # noqa: BLE001
                fallback_text = ""

    # FR-26 AC5: Chart/SmartArt GraphicFrame shapes have has_text_frame==False,
    # so fall back to XML-level extraction (chart part / inline diagram runs).
    if not fallback_text:
        fallback_text = _extract_graphic_frame_fallback_text(shape)

    return OtherShapeIR(
        shape_id=sid,
        name=name,
        kind=ShapeKind.OTHER,
        mso_shape_type=mso_label,
        fallback_text=fallback_text,
    )


def _graphic_frame_type_label(shape: object) -> str | None:
    """Best-effort MSO type label for a GraphicFrame shape (FR-26 AC6).

    python-pptx's ``GraphicFrame.shape_type`` returns None for content it
    does not specifically recognise (notably SmartArt/diagram), leaving
    ``_parse_other`` no label to record. This inspects the shape's own
    ``a:graphicData/@uri`` to distinguish a SmartArt diagram from any other
    unrecognised graphic-frame content.

    Returns None if the shape has no graphicData (i.e. is not a
    GraphicFrame-like element) or the uri cannot be read; the caller falls
    back to "UNKNOWN" in that case.
    """
    element = getattr(shape, "element", None)
    if element is None:
        return None
    try:
        graphic_data = element.find(f".//{_A_GRAPHICDATA_QN}")
    except Exception:  # noqa: BLE001
        return None
    if graphic_data is None:
        return None
    uri = str(graphic_data.get("uri", "") or "")
    if uri == _DIAGRAM_GRAPHIC_DATA_URI:
        return "DIAGRAM"
    if uri:
        return "GRAPHIC_FRAME"
    return None


def _extract_graphic_frame_fallback_text(shape: object) -> str:
    """Best-effort text extraction for Chart/SmartArt GraphicFrame shapes (FR-26 AC5).

    Walks the shape's own XML subtree for ``<a:t>`` runs, which covers inline
    text such as simplified SmartArt fixtures. For charts, python-pptx stores
    the actual chart XML in a related part (not inline in the graphicFrame
    element itself), so the linked chart's XML is also walked for ``<a:t>``
    (titles, data labels) and cached category/series strings (``<c:v>``
    inside ``<c:strCache>``).

    Any failure (missing relationship, malformed XML, etc.) is absorbed and
    "" is returned — this function never raises (ADR-204).
    """
    texts: list[str] = []

    def _collect_a_t(elem: Any) -> None:
        try:
            for t in elem.iter(_A_T_QN):
                if t.text and t.text.strip():
                    texts.append(t.text.strip())
        except Exception:  # noqa: BLE001
            pass

    element = getattr(shape, "element", None)
    if element is not None:
        _collect_a_t(element)

    if getattr(shape, "has_chart", False):
        chart_element: Any = None
        try:
            chart_element = getattr(shape, "chart").element
        except Exception:  # noqa: BLE001
            chart_element = None
        if chart_element is not None:
            _collect_a_t(chart_element)
            try:
                for strcache in chart_element.iter(_C_STRCACHE_QN):
                    for v in strcache.iter(_C_V_QN):
                        if v.text and v.text.strip():
                            texts.append(v.text.strip())
            except Exception:  # noqa: BLE001
                pass

    seen: set[str] = set()
    unique: list[str] = []
    for t in texts:
        if t not in seen:
            seen.add(t)
            unique.append(t)
    return " ".join(unique)


# ---------------------------------------------------------------------------
# Slide metadata helpers
# ---------------------------------------------------------------------------


def _safe_title(slide: object) -> str:
    """Extract title text from a slide's title placeholder (ADR-202).

    Returns "" if the placeholder is absent or empty.
    """
    shapes = getattr(slide, "shapes", None)
    if shapes is None:
        return ""
    title_shape = getattr(shapes, "title", None)
    if title_shape is None:
        return ""
    try:
        text_frame = getattr(title_shape, "text_frame", None)
        if text_frame is not None:
            return str(getattr(text_frame, "text", "") or "")
    except Exception:  # noqa: BLE001
        pass
    return ""


def _safe_notes(slide: object) -> str:
    """Extract notes text from a slide's notes slide (ADR-202).

    Returns "" if the notes slide is absent or has no text.
    """
    try:
        notes_slide = getattr(slide, "notes_slide", None)
        if notes_slide is None:
            return ""
        notes_tf = getattr(notes_slide, "notes_text_frame", None)
        if notes_tf is None:
            return ""
        return str(getattr(notes_tf, "text", "") or "")
    except Exception:  # noqa: BLE001
        return ""


# ---------------------------------------------------------------------------
# Text extraction helper
# ---------------------------------------------------------------------------


def _extract_paragraph_text(para: object) -> str:
    """Concatenate run texts within a paragraph (§8.3, §5.4).

    Empty paragraphs are preserved (text == "") — the assembler decides
    whether to collapse them (ADR-202/§5.4).
    """
    runs = getattr(para, "runs", [])
    if runs:
        return "".join(str(getattr(run, "text", "") or "") for run in runs)
    # Fallback: para.text (includes run texts joined)
    return str(getattr(para, "text", "") or "")
