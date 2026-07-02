"""Shared pytest fixtures (skill §5: fixtures concentrated here).

M2: in-memory PPTX generation fixtures for parser tests (ADR-207).
python-pptx is NOT mocked — real API calls only.
"""

from __future__ import annotations

import io
import struct
import zlib
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Minimal 1×1 PNG helper (for image fixture)
# ---------------------------------------------------------------------------


def _make_png_bytes() -> bytes:
    """Return a minimal 1x1 white PNG as bytes (no Pillow dependency)."""
    # PNG signature
    sig = b"\x89PNG\r\n\x1a\n"

    def _chunk(name: bytes, data: bytes) -> bytes:
        length = struct.pack(">I", len(data))
        crc = struct.pack(">I", zlib.crc32(name + data) & 0xFFFFFFFF)
        return length + name + data + crc

    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)  # 1x1 RGB
    ihdr = _chunk(b"IHDR", ihdr_data)

    # IDAT: filtered scanline for 1x1 RGB (filter byte 0 + RGB white)
    raw = b"\x00\xff\xff\xff"
    compressed = zlib.compress(raw)
    idat = _chunk(b"IDAT", compressed)

    iend = _chunk(b"IEND", b"")
    return sig + ihdr + idat + iend


_PNG_BYTES = _make_png_bytes()


# ---------------------------------------------------------------------------
# In-memory PPTX factories
# ---------------------------------------------------------------------------


def _save_to_bytes(prs: PptxPresentation) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_with_text_slide() -> bytes:
    """Single slide with a title placeholder and a body text box.

    Used for: FR-03 AC1/AC2/AC4, FR-04 AC1.
    """
    prs = PptxPresentation()
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)

    # Set title
    title_shape = slide.shapes.title
    if title_shape is not None:
        title_shape.text = "Test Slide Title"

    # Set body content
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "First paragraph"
    p2 = tf.add_paragraph()
    p2.text = "Second paragraph"
    p2.level = 1

    return _save_to_bytes(prs)


@pytest.fixture
def pptx_with_notes() -> bytes:
    """Single slide with notes text.

    Used for: FR-03 AC3.
    """
    prs = PptxPresentation()
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = "Speaker notes text"

    return _save_to_bytes(prs)


@pytest.fixture
def pptx_multi_slide() -> bytes:
    """Three slides in order: titled, blank, titled again.

    Used for: FR-03 AC1/AC4 (order, index).
    """
    prs = PptxPresentation()

    for i in range(3):
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        title_shape = slide.shapes.title
        if title_shape is not None:
            title_shape.text = f"Slide {i}"

    return _save_to_bytes(prs)


@pytest.fixture
def pptx_no_title_slide() -> bytes:
    """Single slide using a blank layout (no title placeholder).

    Used for: FR-03 AC5 (title absent -> "").
    """
    prs = PptxPresentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank — no title placeholder
    return _save_to_bytes(prs)


@pytest.fixture
def pptx_zero_slides() -> bytes:
    """Presentation with 0 slides.

    Used for: FR-03 AC7.
    """
    prs = PptxPresentation()
    # Do not add any slides — python-pptx starts with 0 slides
    return _save_to_bytes(prs)


@pytest.fixture
def pptx_with_table() -> bytes:
    """Single slide with a 2x3 table.

    Used for: FR-04 AC2/AC7.
    """
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    tbl = slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(6), Inches(2)).table
    tbl.cell(0, 0).text = "A"
    tbl.cell(0, 1).text = "B"
    tbl.cell(0, 2).text = "C"
    tbl.cell(1, 0).text = "D"
    tbl.cell(1, 1).text = "E"
    tbl.cell(1, 2).text = "F"

    return _save_to_bytes(prs)


@pytest.fixture
def pptx_with_empty_table() -> bytes:
    """Single slide with a table where all cells are empty.

    Used for: FR-04 AC7 (empty cells -> "").
    """
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    # Cells remain empty (default)

    return _save_to_bytes(prs)


@pytest.fixture
def pptx_with_image(tmp_path: Path) -> bytes:
    """Single slide with a PNG image shape.

    Used for: FR-04 AC3.
    """
    # Write the minimal PNG to a temp file so python-pptx can add it
    png_file = tmp_path / "test.png"
    png_file.write_bytes(_PNG_BYTES)

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(png_file), Inches(1), Inches(1), Inches(3), Inches(2))

    return _save_to_bytes(prs)


@pytest.fixture
def pptx_with_group() -> bytes:
    """Single slide with a group shape containing one text child.

    Uses XML direct assembly because python-pptx 1.0.2 does not expose
    group_shapes().  The resulting PPTX is valid and python-pptx can
    read it back with shape_type == MSO_SHAPE_TYPE.GROUP.

    Used for: FR-04 AC4.
    """
    from lxml import etree  # noqa: PLC0415 — local import keeps top-level clean

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    sp_tree = slide.shapes._spTree

    # Manually build a grpSp element with one text child.
    grp_xml = (
        "<p:grpSp"
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        "<p:nvGrpSpPr>"
        '<p:cNvPr id="10" name="Group 1"/>'
        "<p:cNvGrpSpPr/>"
        "<p:nvPr/>"
        "</p:nvGrpSpPr>"
        "<p:grpSpPr>"
        "<a:xfrm>"
        '<a:off x="0" y="0"/><a:ext cx="6858000" cy="4572000"/>'
        '<a:chOff x="0" y="0"/><a:chExt cx="6858000" cy="4572000"/>'
        "</a:xfrm>"
        "</p:grpSpPr>"
        "<p:sp>"
        "<p:nvSpPr>"
        '<p:cNvPr id="11" name="TextInGroup"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        "<p:nvPr/>"
        "</p:nvSpPr>"
        "<p:spPr>"
        "<a:xfrm>"
        '<a:off x="914400" y="914400"/>'
        '<a:ext cx="2743200" cy="914400"/>'
        "</a:xfrm>"
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        "</p:spPr>"
        "<p:txBody>"
        "<a:bodyPr/><a:lstStyle/>"
        "<a:p><a:r><a:t>child text</a:t></a:r></a:p>"
        "</p:txBody>"
        "</p:sp>"
        "</p:grpSp>"
    )
    grp_elem = etree.fromstring(grp_xml)
    sp_tree.append(grp_elem)

    return _save_to_bytes(prs)


@pytest.fixture
def pptx_mixed_shapes(tmp_path: Path) -> bytes:
    """Single slide with text + table + image (mixed types, z-order test).

    Used for: FR-04 AC5.
    """
    png_file = tmp_path / "test.png"
    png_file.write_bytes(_PNG_BYTES)

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add in z-order: text, table, image
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.8))
    tb.text_frame.text = "Text shape"

    slide.shapes.add_table(2, 2, Inches(0.5), Inches(1.5), Inches(3), Inches(1.5))

    slide.shapes.add_picture(
        str(png_file), Inches(4), Inches(1), Inches(2), Inches(1.5)
    )

    return _save_to_bytes(prs)


@pytest.fixture
def other_shape_pptx_path() -> Path:
    """Path to tests/fixtures/other_shape.pptx (chart = OTHER shape).

    Used for: FR-04 AC6.
    Created by the fixture generation script; checked into repo (ADR-207).
    """
    return Path(__file__).parent / "fixtures" / "other_shape.pptx"


# ---------------------------------------------------------------------------
# FR-26 (#67): Parser/GraphicFrame coverage fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def pptx_group_three_children(tmp_path: Path) -> bytes:
    """Single slide with one group containing text + table + image children.

    Built via python-pptx's public ``add_group_shape()`` API (no raw XML
    needed for this shape mix). Used for: FR-26 AC1.
    """
    png_file = tmp_path / "group_child.png"
    png_file.write_bytes(_PNG_BYTES)

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1))
    tb.text_frame.text = "group text child"
    tbl_shape = slide.shapes.add_table(
        2, 2, Inches(0.5), Inches(2), Inches(3), Inches(1.5)
    )
    pic = slide.shapes.add_picture(
        str(png_file), Inches(4), Inches(0.5), Inches(2), Inches(1.5)
    )

    slide.shapes.add_group_shape([tb, tbl_shape, pic])

    return _save_to_bytes(prs)


@pytest.fixture
def pptx_nested_group_two_levels() -> bytes:
    """Single slide with a 2-level nested group: outer > inner > leaf text.

    Outer group also has a sibling leaf text shape, so the tree has
    5 total ShapeIR nodes (outer group, inner group, leaf text in inner
    group, sibling text, and the outer group's own node — see test for the
    exact expected iter_shapes() count). Built via ``add_group_shape()``.

    Used for: FR-26 AC2.
    """
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    leaf = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1))
    leaf.text_frame.text = "deep leaf text"
    sibling = slide.shapes.add_textbox(Inches(3), Inches(0.5), Inches(2), Inches(1))
    sibling.text_frame.text = "outer sibling text"

    inner = slide.shapes.add_group_shape([leaf])
    slide.shapes.add_group_shape([inner, sibling])

    return _save_to_bytes(prs)


@pytest.fixture
def pptx_chart_with_title() -> bytes:
    """Single slide with a bar chart that has a chart title set.

    Chart title text lives in the linked chart part's XML (``<a:t>``), not
    inline in the slide's ``p:graphicFrame`` element — exercises the
    chart-part-aware branch of fallback text extraction.

    Used for: FR-26 AC5 (positive extraction), AC6 (mso_shape_type=="CHART").
    """
    from pptx.chart.data import CategoryChartData  # noqa: PLC0415
    from pptx.enum.chart import XL_CHART_TYPE  # noqa: PLC0415

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3"]
    chart_data.add_series("Revenue", (10, 20, 30))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(3),
        chart_data,
    )
    chart = gframe.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Sales Chart"

    return _save_to_bytes(prs)


@pytest.fixture
def pptx_smartart_stub() -> bytes:
    """Single slide with a simplified SmartArt-like GraphicFrame.

    Real SmartArt stores its node text in a separate diagram-data part
    reached via a relationship, which is expensive to assemble faithfully
    for a unit-test fixture. Per FR-26 scope, this fixture instead embeds
    ``<a:t>`` runs directly inside the ``a:graphicData`` element (uri =
    the DrawingML diagram namespace) as a stand-in for those node captions,
    to exercise the inline-XML fallback text extraction path without
    introducing a real diagram-data relationship.

    Used for: FR-26 AC5 (non-chart extraction path), AC6
    (mso_shape_type=="DIAGRAM").
    """
    from lxml import etree  # noqa: PLC0415

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sp_tree = slide.shapes._spTree

    xml = (
        "<p:graphicFrame"
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        "<p:nvGraphicFramePr>"
        '<p:cNvPr id="90" name="SmartArt Stub"/>'
        "<p:cNvGraphicFramePr/>"
        "<p:nvPr/>"
        "</p:nvGraphicFramePr>"
        "<p:xfrm>"
        '<a:off x="914400" y="914400"/><a:ext cx="1828800" cy="1828800"/>'
        "</p:xfrm>"
        "<a:graphic>"
        '<a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
        "<a:p><a:r><a:t>Node Alpha</a:t></a:r></a:p>"
        "<a:p><a:r><a:t>Node Beta</a:t></a:r></a:p>"
        "</a:graphicData>"
        "</a:graphic>"
        "</p:graphicFrame>"
    )
    sp_tree.append(etree.fromstring(xml))

    return _save_to_bytes(prs)


@pytest.fixture
def pptx_images_including_grouped(tmp_path: Path) -> bytes:
    """Single slide with 3 images total: 2 standalone + 1 inside a group.

    Used for: FR-26 AC8 (no image loss across group nesting).
    """
    png_file = tmp_path / "golden.png"
    png_file.write_bytes(_PNG_BYTES)

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.shapes.add_picture(
        str(png_file), Inches(0.5), Inches(0.5), Inches(2), Inches(1.5)
    )
    slide.shapes.add_picture(
        str(png_file), Inches(3), Inches(0.5), Inches(2), Inches(1.5)
    )
    grouped_pic = slide.shapes.add_picture(
        str(png_file), Inches(0.5), Inches(3), Inches(2), Inches(1.5)
    )
    slide.shapes.add_group_shape([grouped_pic])

    return _save_to_bytes(prs)
