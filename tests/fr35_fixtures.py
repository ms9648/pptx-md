"""Synthetic PPTX fixture builder for the FR-35 regression gate (issue #96, W-G).

NOT a test module (no ``test_*`` functions, no pytest collection) — a shared
helper imported by ``tests/test_fr35_regression_gate.py``. Kept in its own
file so the "routing golden" (AC2) and "mock provider combined golden" (AC3)
tests share a single deterministic deck definition.

Provenance (FR-35 AC7 / REQ-v020 §7 픽스처 정책): every shape below is built
**in-process** with ``python-pptx`` — no in-house/customer file of any kind
is read anywhere in this module (no disk reads at all besides the caller
writing the returned bytes out). Safe for public CI / PyPI distribution.

The 4-slide deck reproduces the three real-world structural failure modes
named in REQ-v020 (S5 4열 매핑 / S7 AS-IS↔TO-BE / S12 KPI+도넛) as minimal
synthetic analogues, plus one pure-text slide as a true-negative control:

    index 0 — S5 analogue: title + 4 side-by-side text boxes on one row
              (multi-column signal, ADR-624 AC2).
    index 1 — S7 analogue: title + 3 boxes + 2 straight connectors
              (connector/flow signal, ADR-624 AC3).
    index 2 — S12 analogue: title + 1 pie chart with a chart title
              (chart/SmartArt signal, ADR-624 AC3).
    index 3 — pure text: title + a single text-frame with two linear
              bullets (AC4 정상-부정 — must stay on the text path).

Verified once (scratchpad probe, not re-asserted here — the routing test
in test_fr35_regression_gate.py is the actual regression assertion) against
the real ``pptx_md.parser.parse_presentation`` + ``pptx_md.complexity``
pipeline on 2026-07-11:
    index 0: complexity_score=4  is_visually_complex=True
    index 1: complexity_score=8  is_visually_complex=True
    index 2: complexity_score=4  is_visually_complex=True
    index 3: complexity_score=0  is_visually_complex=False
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

from pptx import Presentation as PptxPresentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches

__all__ = [
    "S5_MULTI_COLUMN_INDEX",
    "S7_CONNECTOR_FLOW_INDEX",
    "S12_CHART_INDEX",
    "PURE_TEXT_INDEX",
    "EXPECTED_COMPLEX_INDICES",
    "EXPECTED_TEXT_PATH_INDICES",
    "EXPECTED_COMPLEXITY_SCORES",
    "build_fr35_synthetic_deck",
    "build_fr35_synthetic_deck_bytes",
]

# ---------------------------------------------------------------------------
# Golden routing constants (AC2) — reviewed, hand-fixed expectations.
# A deliberate change to complexity.py's calibration (⟨DEC-2⟩) that shifts
# these must update this golden through a reviewed PR (FR-35 AC5, D-6
# hard-fail precedent) — it is not expected to drift silently.
# ---------------------------------------------------------------------------
S5_MULTI_COLUMN_INDEX: int = 0
S7_CONNECTOR_FLOW_INDEX: int = 1
S12_CHART_INDEX: int = 2
PURE_TEXT_INDEX: int = 3

EXPECTED_COMPLEX_INDICES: frozenset[int] = frozenset(
    {S5_MULTI_COLUMN_INDEX, S7_CONNECTOR_FLOW_INDEX, S12_CHART_INDEX}
)
EXPECTED_TEXT_PATH_INDICES: frozenset[int] = frozenset({PURE_TEXT_INDEX})

EXPECTED_COMPLEXITY_SCORES: dict[int, int] = {
    S5_MULTI_COLUMN_INDEX: 4,
    S7_CONNECTOR_FLOW_INDEX: 8,
    S12_CHART_INDEX: 4,
    PURE_TEXT_INDEX: 0,
}


def build_fr35_synthetic_deck(path: Path) -> None:
    """Write the 4-slide synthetic S5/S7/S12/text deck to *path* (a .pptx file).

    Pure ``python-pptx`` construction — no external/in-house file is read
    (FR-35 AC7).
    """
    path.write_bytes(build_fr35_synthetic_deck_bytes())


def build_fr35_synthetic_deck_bytes() -> bytes:
    """Return the 4-slide synthetic deck as in-memory .pptx bytes (FR-35 AC7)."""
    prs = PptxPresentation()
    blank = prs.slide_layouts[6]

    # --- index 0: S5 analogue — multi-column mapping row ------------------
    s5 = prs.slides.add_slide(blank)
    _textbox(s5, "4열 매핑", left=0.3, top=0.2, width=9, height=0.6)
    for i, label in enumerate(["기획", "설계", "개발", "운영"]):
        _textbox(s5, label, left=0.3 + i * 2.2, top=1.5, width=2.0, height=1.0)

    # --- index 1: S7 analogue — AS-IS -> TO-BE connector flow --------------
    s7 = prs.slides.add_slide(blank)
    _textbox(s7, "AS-IS -> TO-BE", left=0.3, top=0.2, width=9, height=0.6)
    _textbox(s7, "AS-IS", left=0.3, top=2, width=2, height=1)
    _textbox(s7, "중간 처리", left=4, top=2, width=2, height=1)
    _textbox(s7, "TO-BE", left=7.5, top=2, width=2, height=1)
    s7.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(2.3), Inches(2.5), Inches(4), Inches(2.5)
    )
    s7.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(6), Inches(2.5), Inches(7.5), Inches(2.5)
    )

    # --- index 2: S12 analogue — KPI/도넛(파이) 차트 ------------------------
    s12 = prs.slides.add_slide(blank)
    _textbox(s12, "KPI 대시보드", left=0.3, top=0.2, width=9, height=0.6)
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("Share", (30, 40, 30))
    gframe = s12.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(1), Inches(1.5), Inches(4), Inches(3), chart_data
    )
    gframe.chart.has_title = True
    gframe.chart.chart_title.text_frame.text = "카테고리 비중"

    # --- index 3: pure text — title + single linear-bullet body -----------
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    if s3.shapes.title is not None:
        s3.shapes.title.text = "요약"
    body_ph = None
    for shape in s3.placeholders:
        if shape.placeholder_format.idx == 1:
            body_ph = shape
            break
    if body_ph is not None:
        tf = body_ph.text_frame
        tf.text = "첫째 항목"
        tf.add_paragraph().text = "둘째 항목"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _textbox(
    slide: Any,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.text_frame.text = text
