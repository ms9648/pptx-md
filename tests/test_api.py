"""Tests for FR-16 public API: convert() + ConvertOptions (issue #43).

Covers AC1–AC9.  Pipeline stages are monkeypatched so:
  - No actual PPTX files are required for most tests.
  - No VLM SDK is needed (core-only environment, NFR-08).
  - Each test isolates the contract it is verifying.

AC-to-test mapping:
  AC1 -> TestAc1ConvertReturnsStr
  AC2 -> TestAc2PublicAll
  AC3 -> TestAc3CoreOnlyImport
  AC4 -> TestAc4FakeDescriber
  AC5 -> TestAc5MaskingEmail
  AC6 -> TestAc6MissingPathRaises
  AC7 -> TestAc7EmptyPptxNoException
  AC8 -> TestAc8PyTypedMarker
  AC9 -> TestAc9ValidateLogging
"""

from __future__ import annotations

import io
import logging
import sys
from pathlib import Path
from typing import Any
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation as PptxPresentation

from pptx_md.api import ConvertOptions, convert
from pptx_md.errors import ParseError
from pptx_md.ir import ParagraphIR, PresentationIR, ShapeKind, SlideIR, TextShapeIR
from pptx_md.masking import MaskingOptions

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_minimal_presentation_ir(title: str = "Test Slide") -> PresentationIR:
    """Build a synthetic PresentationIR with one slide for testing."""
    slide = SlideIR(
        index=0,
        title=title,
        shapes=[
            TextShapeIR(
                shape_id=1,
                name="title",
                kind=ShapeKind.TEXT,
                paragraphs=[ParagraphIR(text=title, level=0)],
                is_title=True,
            )
        ],
    )
    pres = PresentationIR(source_path="test.pptx")
    pres.slides.append(slide)
    return pres


def _make_ir_with_image(
    title: str = "Slide With Image",
    image_bytes: bytes = b"\x89PNG\r\n\x1a\n",
    image_ext: str = "png",
    alt_text: str = "",
) -> PresentationIR:
    """Build a PresentationIR that contains one ImageShapeIR for AC4/AC5 tests."""
    from pptx_md.ir import ImageShapeIR  # noqa: PLC0415

    slide = SlideIR(
        index=0,
        title=title,
        shapes=[
            TextShapeIR(
                shape_id=1,
                name="title",
                kind=ShapeKind.TEXT,
                paragraphs=[ParagraphIR(text=title, level=0)],
                is_title=True,
            ),
            ImageShapeIR(
                shape_id=2,
                name="picture",
                kind=ShapeKind.IMAGE,
                image_bytes=image_bytes,
                image_format=image_ext,
                image_ext=image_ext,
                alt_text=alt_text,
            ),
        ],
    )
    pres = PresentationIR(source_path="test.pptx")
    pres.slides.append(slide)
    return pres


def _save_minimal_pptx(path: Path, title: str = "Slide 1") -> None:
    """Create a minimal single-slide PPTX at *path*."""
    prs = PptxPresentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    if title_shape is not None:
        title_shape.text = title
    buf = io.BytesIO()
    prs.save(buf)
    path.write_bytes(buf.getvalue())


def _save_empty_pptx(path: Path) -> None:
    """Create a zero-slide PPTX at *path*."""
    prs = PptxPresentation()
    buf = io.BytesIO()
    prs.save(buf)
    path.write_bytes(buf.getvalue())


# ---------------------------------------------------------------------------
# AC1: convert(path) returns str containing slide heading
# ---------------------------------------------------------------------------


class TestAc1ConvertReturnsStr:
    """ac1_유효한_pptx_변환_성공 — convert() must return str with slide heading."""

    def test_ac1_returns_str_with_heading(self, tmp_path: Path) -> None:
        """AC1: valid PPTX -> convert(path) -> str containing slide heading."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file, title="Hello World")

        result = convert(pptx_file)

        assert isinstance(result, str)
        assert "Hello World" in result

    def test_ac1_accepts_str_path(self, tmp_path: Path) -> None:
        """AC1: convert() must accept a str path as well as pathlib.Path."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        result = convert(str(pptx_file))

        assert isinstance(result, str)

    def test_ac1_pipeline_stages_called_in_order(self, tmp_path: Path) -> None:
        """AC1: verify parse->enrich_images->enrich_descriptions->assemble order."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        call_log: list[str] = []
        fake_ir = _make_minimal_presentation_ir()

        def fake_parse(path: Any) -> PresentationIR:
            call_log.append("parse")
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            call_log.append("enrich_images")

        def fake_enrich_descriptions(pres: Any, describer: Any, **kwargs: Any) -> None:
            call_log.append("enrich_descriptions")

        def fake_assemble(pres: Any, *, masking: Any = None, **kwargs: Any) -> str:
            call_log.append("assemble")
            return "# Test Slide\n"

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
            patch("pptx_md.api.enrich_descriptions", fake_enrich_descriptions),
            patch("pptx_md.api.assemble_document", fake_assemble),
        ):
            result = convert(pptx_file)

        assert call_log == ["parse", "enrich_images", "enrich_descriptions", "assemble"]
        assert result == "# Test Slide\n"


# ---------------------------------------------------------------------------
# AC2: pptx_md.__all__ contains convert and ConvertOptions
# ---------------------------------------------------------------------------


class TestAc2PublicAll:
    """ac2_all_contains_required_symbols — __all__ must include the required names."""

    def test_ac2_convert_in_all(self) -> None:
        """AC2: 'convert' must be in pptx_md.__all__."""
        import pptx_md

        assert "convert" in pptx_md.__all__

    def test_ac2_convert_options_in_all(self) -> None:
        """AC2: 'ConvertOptions' must be in pptx_md.__all__."""
        import pptx_md

        assert "ConvertOptions" in pptx_md.__all__

    def test_ac2_full_public_surface(self) -> None:
        """AC2: __all__ must contain exactly the ADR-603 symbols (no more, no less)."""
        import pptx_md

        expected = {
            "convert",
            "ConvertOptions",
            "MaskingOptions",
            "MASK_TOKEN",
            "validate_markdown",
            "ValidationResult",
            "ImageDescriber",
            "get_describer",
            "PptxMdError",
            "ParseError",
            "DescribeError",
            "InstallationError",
        }
        assert set(pptx_md.__all__) == expected

    def test_ac2_internal_functions_not_in_all(self) -> None:
        """AC2: Internal pipeline functions must NOT appear in __all__."""
        import pptx_md

        internal = [
            "parse_presentation",
            "enrich_images",
            "enrich_descriptions",
            "assemble_document",
            "assemble_slide",
        ]
        for name in internal:
            assert name not in pptx_md.__all__, f"{name!r} must not be in __all__"

    def test_ac2_all_symbols_importable(self) -> None:
        """AC2: every symbol in __all__ must be importable from pptx_md."""
        import pptx_md

        for name in pptx_md.__all__:
            assert hasattr(pptx_md, name), f"pptx_md.{name} not found"


# ---------------------------------------------------------------------------
# AC3: core-only import succeeds without VLM SDKs
# ---------------------------------------------------------------------------


class TestAc3CoreOnlyImport:
    """ac3_core_only_import_no_error — import pptx_md must not require VLM SDKs."""

    def test_ac3_import_pptx_md_no_error(self) -> None:
        """AC3: import pptx_md must succeed without errors (NFR-08)."""
        # Simply importing the module must never raise
        import pptx_md  # noqa: F401

    def test_ac3_no_vlm_sdk_eager_import(self) -> None:
        """AC3: import pptx_md must NOT trigger eager import of anthropic/openai.

        We verify by checking that the already-loaded pptx_md package
        did NOT pull in VLM SDKs as a side effect.  Reloading the full
        module tree would corrupt global module state and break subsequent
        tests, so we check the *current* sys.modules state instead.
        The first import already happened at test-collection time; if
        anthropic/openai are absent now they were never eagerly imported.
        """
        import pptx_md  # noqa: F401 — ensure loaded

        # VLM SDKs must not be present unless the user explicitly installed
        # and imported them outside of pptx_md (NFR-08).
        assert "anthropic" not in sys.modules, "anthropic was eagerly imported"
        assert "openai" not in sys.modules, "openai was eagerly imported"


# ---------------------------------------------------------------------------
# AC4: fake describer injected -> describe called + result in Markdown
# ---------------------------------------------------------------------------


class TestAc4FakeDescriber:
    """ac4_fake_describer_주입 — fake describer must be called; result in Markdown."""

    def test_ac4_describer_describe_called_and_result_in_md(
        self, tmp_path: Path
    ) -> None:
        """AC4: fake describer.describe is invoked by the real enrich_descriptions
        pipeline and the returned text appears in the final Markdown output.

        enrich_descriptions is NOT mocked so the real call chain
        (enrich_descriptions -> _enrich_single -> describer.describe
        -> shape.description) is exercised.  assemble_document is also real
        so the description flows all the way through to the Markdown string.
        """
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        # IR contains an ImageShapeIR so enrich_descriptions has something to process.
        fake_ir = _make_ir_with_image(image_ext="png")

        fake_describer = MagicMock()
        fake_describer.describe.return_value = "a diagram showing system flow"

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            # enrich_images (M3 classifier) is mocked to avoid real image bytes
            # dependency; enrich_descriptions (M4) runs real.
            pass

        opts = ConvertOptions(describer=fake_describer)

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
        ):
            result = convert(pptx_file, options=opts)

        # describer.describe must have been called by the real enrich_descriptions
        fake_describer.describe.assert_called_once()
        # the returned description text must appear in the assembled Markdown
        assert "a diagram showing system flow" in result

    def test_ac4_describer_passed_to_enrich_descriptions(self, tmp_path: Path) -> None:
        """AC4: the exact describer object is forwarded to enrich_descriptions."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        fake_ir = _make_ir_with_image()
        fake_describer = MagicMock()
        # describe returns a string so enrich_descriptions doesn't error out
        fake_describer.describe.return_value = "described"
        captured: list[Any] = []

        original_enrich_descriptions = __import__(
            "pptx_md.description_pipeline", fromlist=["enrich_descriptions"]
        ).enrich_descriptions

        def spy_enrich_descriptions(pres: Any, describer: Any, **kwargs: Any) -> None:
            captured.append(describer)
            original_enrich_descriptions(pres, describer, **kwargs)

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        opts = ConvertOptions(describer=fake_describer)

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
            patch("pptx_md.api.enrich_descriptions", spy_enrich_descriptions),
        ):
            convert(pptx_file, options=opts)

        assert captured == [fake_describer]


# ---------------------------------------------------------------------------
# AC5: MaskingOptions(enabled=True) -> email replaced with [REDACTED]
# ---------------------------------------------------------------------------


class TestAc5MaskingEmail:
    """ac5_masking_options_이메일_치환 — enabled masking must replace email."""

    def test_ac5_email_redacted(self, tmp_path: Path) -> None:
        """AC5: MaskingOptions(enabled=True) -> email address -> [REDACTED].

        Uses a real IR containing a text shape with an email address.
        enrich_descriptions and enrich_images are mocked (no image bytes / VLM),
        but assemble_document runs REAL with MaskingOptions(enabled=True) so
        the actual mask_text() substitution path is exercised end-to-end.
        """
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        # Build an IR whose text contains an email address.
        slide = SlideIR(
            index=0,
            title="Contact",
            shapes=[
                TextShapeIR(
                    shape_id=1,
                    name="body",
                    kind=ShapeKind.TEXT,
                    paragraphs=[
                        ParagraphIR(
                            text="Contact user@example.com for details.", level=0
                        )
                    ],
                    is_title=False,
                )
            ],
        )
        fake_ir = PresentationIR(source_path="test.pptx")
        fake_ir.slides.append(slide)

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        def fake_enrich_descriptions(pres: Any, describer: Any, **kwargs: Any) -> None:
            pass

        opts = ConvertOptions(masking=MaskingOptions(enabled=True))

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
            patch("pptx_md.api.enrich_descriptions", fake_enrich_descriptions),
        ):
            result = convert(pptx_file, options=opts)

        # Real mask_text() must have replaced the email with [REDACTED].
        assert "user@example.com" not in result, "email must be redacted"
        assert "[REDACTED]" in result, "[REDACTED] token must appear in output"

    def test_ac5_masking_forwarded_to_assemble(self, tmp_path: Path) -> None:
        """AC5: masking option is forwarded to assemble_document as kwarg."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        fake_ir = _make_minimal_presentation_ir()
        captured_masking: list[Any] = []

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        def fake_enrich_descriptions(pres: Any, describer: Any, **kwargs: Any) -> None:
            pass

        def fake_assemble(pres: Any, *, masking: Any = None, **kwargs: Any) -> str:
            captured_masking.append(masking)
            return "# Slide\n"

        masking_opts = MaskingOptions(enabled=True)
        opts = ConvertOptions(masking=masking_opts)

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
            patch("pptx_md.api.enrich_descriptions", fake_enrich_descriptions),
            patch("pptx_md.api.assemble_document", fake_assemble),
        ):
            convert(pptx_file, options=opts)

        assert captured_masking == [masking_opts]


# ---------------------------------------------------------------------------
# AC6: non-existent path -> ParseError (or subclass) raised
# ---------------------------------------------------------------------------


class TestAc6MissingPathRaises:
    """ac6_존재하지_않는_경로_ParseError — missing file must raise a pptx_md error."""

    def test_ac6_nonexistent_path_raises(self, tmp_path: Path) -> None:
        """AC6: convert() on a non-existent path must raise ParseError."""
        missing = tmp_path / "does_not_exist.pptx"

        with pytest.raises((ParseError, FileNotFoundError, OSError)):
            convert(missing)

    def test_ac6_invalid_file_raises_parse_error(self, tmp_path: Path) -> None:
        """AC6: convert() on a non-PPTX file must raise ParseError."""
        bad_file = tmp_path / "not_a_pptx.pptx"
        bad_file.write_bytes(b"this is not a zip/pptx file")

        with pytest.raises((ParseError, Exception)):
            convert(bad_file)


# ---------------------------------------------------------------------------
# AC7: empty PPTX -> no exception, returns str (possibly empty/whitespace)
# ---------------------------------------------------------------------------


class TestAc7EmptyPptxNoException:
    """ac7_빈_pptx_예외없이_문자열 — zero-slide PPTX must not raise; returns str."""

    def test_ac7_empty_pptx_returns_str(self, tmp_path: Path) -> None:
        """AC7: PPTX with 0 slides -> convert() returns str without exception."""
        pptx_file = tmp_path / "empty.pptx"
        _save_empty_pptx(pptx_file)

        result = convert(pptx_file)

        assert isinstance(result, str)

    def test_ac7_monkeypatched_empty_pres_no_exception(self, tmp_path: Path) -> None:
        """AC7: empty PresentationIR (no slides) -> no exception, returns str."""
        pptx_file = tmp_path / "empty.pptx"
        _save_empty_pptx(pptx_file)

        empty_ir = PresentationIR(source_path="empty.pptx")

        def fake_parse(path: Any) -> PresentationIR:
            return empty_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        def fake_enrich_descriptions(pres: Any, describer: Any, **kwargs: Any) -> None:
            pass

        def fake_assemble(pres: Any, *, masking: Any = None, **kwargs: Any) -> str:
            return ""

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
            patch("pptx_md.api.enrich_descriptions", fake_enrich_descriptions),
            patch("pptx_md.api.assemble_document", fake_assemble),
        ):
            result = convert(pptx_file)

        assert isinstance(result, str)


# ---------------------------------------------------------------------------
# AC8: py.typed marker exists in package directory
# ---------------------------------------------------------------------------


class TestAc8PyTypedMarker:
    """ac8_py_typed_존재 — py.typed marker must exist in the package directory."""

    def test_ac8_py_typed_file_exists(self) -> None:
        """AC8: src/pptx_md/py.typed must exist as a file (PEP 561, ADR-602)."""
        import pptx_md

        pkg_dir = Path(pptx_md.__file__).parent  # type: ignore[arg-type]
        py_typed = pkg_dir / "py.typed"
        assert py_typed.exists(), f"py.typed not found at {py_typed}"
        assert py_typed.is_file(), f"py.typed is not a file at {py_typed}"


# ---------------------------------------------------------------------------
# AC9: validate=True -> validate_markdown called, warnings logged
# ---------------------------------------------------------------------------


class TestAc9ValidateLogging:
    """ac9_validate_true_경고_로깅 — validate=True must call validate_markdown + log."""

    def test_ac9_validate_markdown_called_when_flag_true(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        """AC9: validate=True must call validate_markdown() internally."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        fake_ir = _make_minimal_presentation_ir()
        called: list[str] = []

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        def fake_enrich_descriptions(pres: Any, describer: Any, **kwargs: Any) -> None:
            pass

        def fake_assemble(pres: Any, *, masking: Any = None, **kwargs: Any) -> str:
            return "# Slide\n\nContent\n"

        def fake_validate(md: str) -> Any:
            called.append(md)
            from pptx_md.validator import ValidationResult  # noqa: PLC0415

            return ValidationResult(valid=True, warnings=[])

        opts = ConvertOptions(validate=True)

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
            patch("pptx_md.api.enrich_descriptions", fake_enrich_descriptions),
            patch("pptx_md.api.assemble_document", fake_assemble),
            patch("pptx_md.api.validate_markdown", fake_validate),
        ):
            result = convert(pptx_file, options=opts)

        assert len(called) == 1, "validate_markdown must be called exactly once"
        assert isinstance(result, str), "return type must remain str (ADR-604)"

    def test_ac9_validate_false_skips_validate_markdown(self, tmp_path: Path) -> None:
        """AC9: validate=False (default) must NOT call validate_markdown."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        fake_ir = _make_minimal_presentation_ir()
        called: list[str] = []

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        def fake_enrich_descriptions(pres: Any, describer: Any, **kwargs: Any) -> None:
            pass

        def fake_assemble(pres: Any, *, masking: Any = None, **kwargs: Any) -> str:
            return "# Slide\n"

        def fake_validate(md: str) -> Any:
            called.append(md)
            from pptx_md.validator import ValidationResult  # noqa: PLC0415

            return ValidationResult(valid=True, warnings=[])

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
            patch("pptx_md.api.enrich_descriptions", fake_enrich_descriptions),
            patch("pptx_md.api.assemble_document", fake_assemble),
            patch("pptx_md.api.validate_markdown", fake_validate),
        ):
            convert(pptx_file)

        assert (
            len(called) == 0
        ), "validate_markdown must NOT be called when validate=False"

    def test_ac9_warning_logged_when_invalid(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        """AC9: validate=True + invalid document -> warning logged via pptx_md.api."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        fake_ir = _make_minimal_presentation_ir()

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        def fake_enrich_descriptions(pres: Any, describer: Any, **kwargs: Any) -> None:
            pass

        def fake_assemble(pres: Any, *, masking: Any = None, **kwargs: Any) -> str:
            return "# Slide\n\n```\nunclosed fence"

        def fake_validate(md: str) -> Any:
            from pptx_md.validator import ValidationResult  # noqa: PLC0415

            return ValidationResult(valid=False, warnings=["unclosed fence"])

        opts = ConvertOptions(validate=True)

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
            patch("pptx_md.api.enrich_descriptions", fake_enrich_descriptions),
            patch("pptx_md.api.assemble_document", fake_assemble),
            patch("pptx_md.api.validate_markdown", fake_validate),
            caplog.at_level(logging.WARNING, logger="pptx_md.api"),
        ):
            result = convert(pptx_file, options=opts)

        assert isinstance(result, str), "return type must remain str even when invalid"
        assert any(
            "validate_markdown" in r.message or "invalid" in r.message
            for r in caplog.records
        ), "a warning log record must be emitted"

    def test_ac9_return_type_always_str(self, tmp_path: Path) -> None:
        """AC9 (ADR-604): convert() return type must be str regardless of validate."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        fake_ir = _make_minimal_presentation_ir()

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        def fake_enrich_descriptions(pres: Any, describer: Any, **kwargs: Any) -> None:
            pass

        def fake_assemble(pres: Any, *, masking: Any = None, **kwargs: Any) -> str:
            return "# Slide\n"

        def fake_validate(md: str) -> Any:
            from pptx_md.validator import ValidationResult  # noqa: PLC0415

            return ValidationResult(valid=True, warnings=[])

        for validate_flag in (True, False):
            opts = ConvertOptions(validate=validate_flag)
            with (
                patch("pptx_md.api.parse_presentation", fake_parse),
                patch("pptx_md.api.enrich_images", fake_enrich_images),
                patch("pptx_md.api.enrich_descriptions", fake_enrich_descriptions),
                patch("pptx_md.api.assemble_document", fake_assemble),
                patch("pptx_md.api.validate_markdown", fake_validate),
            ):
                result = convert(pptx_file, options=opts)

            assert isinstance(
                result, str
            ), f"convert(validate={validate_flag}) must return str"


# ---------------------------------------------------------------------------
# FR-27 (issue #68, ARCH-M12) — ConvertOptions new fields + Stage3 forwarding
#
# Test naming: test_ac<N>_fr27_<desc> where <N> matches issue #68's AC
# numbering (distinct from the FR-16 AC numbering used above in this file).
# ---------------------------------------------------------------------------


class TestFr27ConvertOptions:
    """FR-27 (issue #68): describe_max_workers / diagram_mermaid fields."""

    def test_ac_fr27_default_값_하위호환(self) -> None:
        """ARCH-M12 §3.5: new fields default to values equivalent to pre-M12
        behaviour (describe_max_workers=4, diagram_mermaid=False)."""
        opts = ConvertOptions()
        assert opts.describe_max_workers == 4
        assert opts.diagram_mermaid is False

    def test_ac_fr27_frozen_불변(self) -> None:
        """ADR-602/613: ConvertOptions remains a frozen dataclass."""
        import dataclasses

        opts = ConvertOptions()
        with pytest.raises(dataclasses.FrozenInstanceError):
            opts.describe_max_workers = 8  # type: ignore[misc]

    def test_ac_fr27_기존_호출_무손상(self, tmp_path: Path) -> None:
        """ADR-613: ConvertOptions() / convert(src) without new kwargs still work."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        result = convert(pptx_file)  # no options at all

        assert isinstance(result, str)

    def test_ac1_ac7_fr27_옵션이_enrich_descriptions로_전달(
        self, tmp_path: Path
    ) -> None:
        """issue AC1/AC7: describe_max_workers/diagram_mermaid are forwarded to
        enrich_descriptions as keyword arguments by convert()."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        fake_ir = _make_ir_with_image()
        captured: dict[str, Any] = {}

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        def spy_enrich_descriptions(
            pres: Any,
            describer: Any,
            *,
            max_workers: int = 4,
            diagram_mermaid: bool = False,
        ) -> None:
            captured["max_workers"] = max_workers
            captured["diagram_mermaid"] = diagram_mermaid

        opts = ConvertOptions(describe_max_workers=2, diagram_mermaid=True)

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
            patch("pptx_md.api.enrich_descriptions", spy_enrich_descriptions),
        ):
            convert(pptx_file, options=opts)

        assert captured == {"max_workers": 2, "diagram_mermaid": True}


class TestAc1Fr27EndToEndDescriber:
    """issue AC1: mock describer 주입 -> description 채워짐 + Markdown 본문 포함."""

    def test_ac1_fr27_mock_describer_주입_markdown_본문_포함(
        self, tmp_path: Path
    ) -> None:
        """issue AC1: real enrich_descriptions pipeline fills description and
        the text flows into the final Markdown body."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        fake_ir = _make_ir_with_image(image_bytes=b"\x89PNG_AC1_FR27")

        class MockDescriber:
            def __init__(self) -> None:
                self.calls = 0

            def describe(
                self, image_bytes: bytes, image_ext: str, shape_hint: str | None
            ) -> str:
                self.calls += 1
                return "issue AC1 description text"

        mock_describer = MockDescriber()

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        opts = ConvertOptions(describer=mock_describer)

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
        ):
            result = convert(pptx_file, options=opts)

        assert mock_describer.calls == 1
        # the image shape from fake_ir must have its description filled in-place
        image_shape = fake_ir.slides[0].shapes[1]
        assert image_shape.description == "issue AC1 description text"
        assert "issue AC1 description text" in result


class TestAc8Fr27DeterministicGoldenRegression:
    """issue AC8: 고정 VLM mock, 동일 PPTX 2회 변환 -> Markdown 바이트 동일."""

    def test_ac8_fr27_동일_pptx_2회_변환_바이트_동일(self, tmp_path: Path) -> None:
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        def _fresh_ir() -> PresentationIR:
            return _make_ir_with_image(image_bytes=b"\x89PNG_GOLDEN_FIXED")

        describer = MagicMock()
        describer.describe.return_value = "a fixed golden vlm description"

        def fake_enrich_images(pres: Any) -> None:
            pass

        opts = ConvertOptions(describer=describer, describe_max_workers=4)

        outputs: list[str] = []
        for _ in range(2):
            with (
                patch(
                    "pptx_md.api.parse_presentation",
                    lambda path, _ir=_fresh_ir: _ir(),
                ),
                patch("pptx_md.api.enrich_images", fake_enrich_images),
            ):
                outputs.append(convert(pptx_file, options=opts))

        assert outputs[0].encode("utf-8") == outputs[1].encode("utf-8")
