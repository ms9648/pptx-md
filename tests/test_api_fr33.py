"""Tests for api.py's FR-33 ConvertOptions fields + Stage 3.5 wiring
(issue #95, ARCH-v020 §3.5, ADR-627).

WBS-O AC list covered here:
    AC3 ConvertOptions 신규 필드 전부 default 안전값(ARCH §3.5 표 그대로).
    AC4 INV-3 바이트 동일: convert(src)/ConvertOptions()(신규 동작 off) ->
        현행 출력과 바이트 동일. 골든/기존 테스트 무손상.
    (겸사) Stage 3.5가 opts.visual_reconstruct=True일 때만 reconstruct_slides에
    올바른 인자로 배선되는지.

Test naming convention: test_ac<N>_<description>.
"""

from __future__ import annotations

import dataclasses
import io
from pathlib import Path
from typing import Any
from unittest.mock import patch

import pytest
from pptx import Presentation as PptxPresentation

from pptx_md.api import ConvertOptions, convert
from pptx_md.ir import ParagraphIR, PresentationIR, ShapeKind, SlideIR, TextShapeIR


def _save_minimal_pptx(path: Path, title: str = "Slide 1") -> None:
    prs = PptxPresentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    if title_shape is not None:
        title_shape.text = title
    buf = io.BytesIO()
    prs.save(buf)
    path.write_bytes(buf.getvalue())


def _make_minimal_presentation_ir(title: str = "Test Slide") -> PresentationIR:
    slide = SlideIR(
        index=0,
        title=title,
        shapes=[
            TextShapeIR(
                shape_id=1,
                name="title",
                kind=ShapeKind.TEXT,
                paragraphs=[ParagraphIR(text=title, level=0)],
                is_title=True,
            )
        ],
    )
    pres = PresentationIR(source_path="test.pptx")
    pres.slides.append(slide)
    return pres


# ===========================================================================
# AC3 — ConvertOptions 신규 필드 전부 default 안전값
# ===========================================================================


class TestAc3ConvertOptionsFr33Fields:
    """ARCH-v020 §3.5 표: 신규 필드 전부 default 안전값."""

    def test_ac3_default_values(self) -> None:
        opts = ConvertOptions()
        assert opts.reconstructor is None
        assert opts.visual_reconstruct is False
        assert opts.max_vlm_slides is None
        assert opts.embed_rendered_image is False
        assert opts.force_render == frozenset()
        assert opts.force_text == frozenset()
        assert opts.render_dpi == 150
        assert opts.reconstruct_max_workers == 4

    def test_ac3_frozen_immutable(self) -> None:
        opts = ConvertOptions()
        with pytest.raises(dataclasses.FrozenInstanceError):
            opts.visual_reconstruct = True  # type: ignore[misc]

    def test_ac3_fields_are_overridable(self) -> None:
        fake_reconstructor = object()
        opts = ConvertOptions(
            reconstructor=fake_reconstructor,  # type: ignore[arg-type]
            visual_reconstruct=True,
            max_vlm_slides=5,
            embed_rendered_image=True,
            force_render=frozenset({1, 2}),
            force_text=frozenset({3}),
            render_dpi=200,
            reconstruct_max_workers=8,
        )
        assert opts.reconstructor is fake_reconstructor
        assert opts.visual_reconstruct is True
        assert opts.max_vlm_slides == 5
        assert opts.embed_rendered_image is True
        assert opts.force_render == frozenset({1, 2})
        assert opts.force_text == frozenset({3})
        assert opts.render_dpi == 200
        assert opts.reconstruct_max_workers == 8


# ===========================================================================
# AC4 (INV-3) — 바이트 동일: convert(src) / ConvertOptions() 무손상
# ===========================================================================


class TestAc4Inv3ByteIdentical:
    """FR-33 AC2/INV-3: 신규 옵션 미지정(off) -> 현행 출력과 바이트 동일."""

    def test_ac4_convert_no_options_unchanged(self, tmp_path: Path) -> None:
        """convert(src) (options=None) 여전히 정상 동작 -- 기존 호출 무손상."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file, title="Hello FR-33")

        result = convert(pptx_file)

        assert isinstance(result, str)
        assert "Hello FR-33" in result

    def test_ac4_convert_options_default_byte_identical_to_no_options(
        self, tmp_path: Path
    ) -> None:
        """ConvertOptions()(전 필드 default) 결과 == convert(src) 결과, 바이트 동일."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file, title="Byte Identical")

        result_implicit = convert(pptx_file)
        result_explicit = convert(pptx_file, options=ConvertOptions())

        assert result_implicit.encode("utf-8") == result_explicit.encode("utf-8")

    def test_ac4_reconstruct_slides_not_invoked_when_visual_reconstruct_false(
        self, tmp_path: Path
    ) -> None:
        """visual_reconstruct=False(기본) -> reconstruct_slides()가 전혀 호출되지
        않는다."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        with patch("pptx_md.api.reconstruct_slides") as mock_reconstruct:
            convert(pptx_file)  # default options
            convert(pptx_file, options=ConvertOptions())  # explicit default

        mock_reconstruct.assert_not_called()

    def test_ac4_output_matches_pre_fr33_rendering(self, tmp_path: Path) -> None:
        """오프 경로 산출물이 shape 렌더 결과를 그대로 담고 있다(회귀 없음)."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file, title="Regression Check")

        fake_ir = _make_minimal_presentation_ir(title="Regression Check")

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        def fake_enrich_descriptions(pres: Any, describer: Any, **kwargs: Any) -> None:
            pass

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
            patch("pptx_md.api.enrich_descriptions", fake_enrich_descriptions),
        ):
            result = convert(pptx_file)

        assert "Regression Check" in result
        assert fake_ir.slides[0].reconstructed_md is None  # untouched by Stage 3.5


# ===========================================================================
# Stage 3.5 wiring — visual_reconstruct=True forwards the right kwargs
# ===========================================================================


class TestStage35Wiring:
    """visual_reconstruct=True일 때 reconstruct_slides()가 올바른 인자로 호출된다."""

    def test_reconstruct_slides_called_with_forwarded_options(
        self, tmp_path: Path
    ) -> None:
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        fake_reconstructor = object()
        opts = ConvertOptions(
            visual_reconstruct=True,
            reconstructor=fake_reconstructor,  # type: ignore[arg-type]
            max_vlm_slides=3,
            force_render=frozenset({0}),
            force_text=frozenset({1}),
            embed_rendered_image=True,
            render_dpi=200,
            reconstruct_max_workers=2,
        )

        with patch("pptx_md.api.reconstruct_slides") as mock_reconstruct:
            convert(pptx_file, options=opts)

        mock_reconstruct.assert_called_once()
        call_args, call_kwargs = mock_reconstruct.call_args
        # positional args: (pres, reconstructor)
        assert call_args[1] is fake_reconstructor
        assert call_kwargs == {
            "max_slides": 3,
            "force_render": frozenset({0}),
            "force_text": frozenset({1}),
            "embed": True,
            "render_dpi": 200,
            "max_workers": 2,
        }

    def test_reconstruct_slides_result_flows_into_markdown(
        self, tmp_path: Path
    ) -> None:
        """reconstruct_slides가 실제로 슬롯을 채우면 그 결과가 최종 Markdown에
        반영된다."""
        pptx_file = tmp_path / "deck.pptx"
        _save_minimal_pptx(pptx_file)

        fake_ir = _make_minimal_presentation_ir(title="Original Title")

        def fake_parse(path: Any) -> PresentationIR:
            return fake_ir

        def fake_enrich_images(pres: Any) -> None:
            pass

        def fake_enrich_descriptions(pres: Any, describer: Any, **kwargs: Any) -> None:
            pass

        def fake_reconstruct(pres: Any, reconstructor: Any, **kwargs: Any) -> None:
            pres.slides[0].reconstructed_md = "## VLM reconstructed body"

        opts = ConvertOptions(visual_reconstruct=True, reconstructor=object())  # type: ignore[arg-type]

        with (
            patch("pptx_md.api.parse_presentation", fake_parse),
            patch("pptx_md.api.enrich_images", fake_enrich_images),
            patch("pptx_md.api.enrich_descriptions", fake_enrich_descriptions),
            patch("pptx_md.api.reconstruct_slides", fake_reconstruct),
        ):
            result = convert(pptx_file, options=opts)

        assert "## VLM reconstructed body" in result
