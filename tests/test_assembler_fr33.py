"""Tests for assembler.py's FR-33 reconstructed_md slot consumption
(issue #95, ARCH-v020 §3.4, ADR-627).

WBS-O AC list covered here:
    AC8 어셈블러: reconstructed_md 존재 시 슬라이드 body로 렌더, 없으면 현행 경로.
        off(None) -> 현행 바이트 동일.
    AC9 (INV-5): assembler.py 최상단 VLM SDK·pypdfium2 import 0 (지연/주입만).

Test naming convention: test_ac<N>_<description>.
"""

from __future__ import annotations

import ast
from pathlib import Path

from pptx_md.assembler import assemble_document, assemble_slide
from pptx_md.ir import (
    ParagraphIR,
    PresentationIR,
    ShapeKind,
    SlideIR,
    TableShapeIR,
    TextShapeIR,
)
from pptx_md.masking import MaskingOptions


def _slide_with_shapes(index: int, title: str = "제목") -> SlideIR:
    return SlideIR(
        index=index,
        title=title,
        shapes=[
            TextShapeIR(
                shape_id=1,
                name="body",
                kind=ShapeKind.TEXT,
                paragraphs=[ParagraphIR(text="original shape body", level=0)],
            ),
            TableShapeIR(
                shape_id=2,
                name="table",
                kind=ShapeKind.TABLE,
                rows=[["a", "b"], ["c", "d"]],
                n_rows=2,
                n_cols=2,
            ),
        ],
    )


# ===========================================================================
# AC8 — reconstructed_md 존재 시 body로 렌더, 원래 shape 렌더는 건너뛴다
# ===========================================================================


def test_ac8_reconstructed_md_replaces_shape_body() -> None:
    """AC8: reconstructed_md가 채워지면 원래 shape 렌더 대신 그 텍스트가 body가 된다."""
    slide = _slide_with_shapes(0)
    slide.reconstructed_md = "## Reconstructed\n\n| x | y |\n| --- | --- |\n| 1 | 2 |"

    rendered = assemble_slide(slide)

    assert "original shape body" not in rendered
    assert "| x | y |" in rendered
    assert "## 제목" in rendered  # heading unaffected by the slot


def test_ac8_heading_still_computed_normally_with_reconstructed_md() -> None:
    """AC8: reconstructed_md 경로에서도 헤딩 계산(제목/fallback)은 그대로 동작한다."""
    slide = SlideIR(index=4, title="", reconstructed_md="body text")
    rendered = assemble_slide(slide)
    assert rendered.startswith("## Slide 5")  # 1-based fallback heading


def test_ac8_document_level_slot_consumption_in_order() -> None:
    """AC8: assemble_document에서도 reconstructed_md 슬라이드와 일반 슬라이드가
    인덱스 순으로 결합된다(ADR-220 병합 순서 계승, FR-33 AC1/AC4)."""
    slide0 = _slide_with_shapes(0, title="A")
    slide1 = _slide_with_shapes(1, title="B")
    slide1.reconstructed_md = "## Reconstructed body for slide 2"
    pres = PresentationIR(
        source_path="doc.pptx", slides=[slide1, slide0]
    )  # 뒤섞인 순서

    doc = assemble_document(pres)

    idx_a = doc.index("## A")
    idx_b = doc.index("## B")
    assert idx_a < idx_b, "슬라이드는 IR 리스트 순서가 아니라 index 순으로 병합된다"
    assert "## Reconstructed body for slide 2" in doc
    assert "original shape body" in doc  # slide0 (off) unaffected


def test_ac8_masking_applied_to_reconstructed_md() -> None:
    """FR-33 AC7 정합: masking이 설정되어 있으면 reconstructed_md에도 적용된다."""
    slide = SlideIR(
        index=0, title="Contact", reconstructed_md="Contact user@example.com now"
    )
    rendered = assemble_slide(slide, masking=MaskingOptions(enabled=True))
    assert "user@example.com" not in rendered
    assert "[REDACTED]" in rendered


# ===========================================================================
# AC8 (INV-3) — reconstructed_md=None(기본) -> 현행 shape 렌더 경로, 바이트 동일
# ===========================================================================


def test_ac8_none_slot_falls_through_to_existing_shape_render() -> None:
    """AC8/INV-3: reconstructed_md=None(기본)이면 기존 shape 렌더 결과와 동일하다."""
    slide_off = _slide_with_shapes(0)
    slide_control = _slide_with_shapes(0)  # 별개 인스턴스, 동일 IR

    rendered_off = assemble_slide(slide_off)
    rendered_control = assemble_slide(slide_control)

    assert slide_off.reconstructed_md is None
    assert rendered_off == rendered_control
    assert "original shape body" in rendered_off


def test_ac8_document_byte_identical_when_all_slots_none() -> None:
    """AC8/INV-3: 전 슬라이드 reconstructed_md=None -> assemble_document 바이트 동일
    (2회 독립 빌드 비교, 결정성)."""
    pres1 = PresentationIR(
        source_path="doc.pptx", slides=[_slide_with_shapes(0), _slide_with_shapes(1)]
    )
    pres2 = PresentationIR(
        source_path="doc.pptx", slides=[_slide_with_shapes(0), _slide_with_shapes(1)]
    )

    doc1 = assemble_document(pres1)
    doc2 = assemble_document(pres2)

    assert doc1 == doc2


# ===========================================================================
# AC9 (INV-5) — assembler.py 최상단 VLM SDK·pypdfium2 import 0
# ===========================================================================


def test_ac9_assembler_no_top_level_sdk_or_pypdfium2_import() -> None:
    """INV-5: assembler.py 최상단에 anthropic/openai/pypdfium2 import가 없다(AST)."""
    src_path = Path(__file__).parent.parent / "src" / "pptx_md" / "assembler.py"
    tree = ast.parse(src_path.read_text(encoding="utf-8"))

    top_level_imports: list[str] = []
    for node in tree.body:
        if isinstance(node, ast.Import):
            for alias in node.names:
                top_level_imports.append(alias.name.split(".")[0])
        elif isinstance(node, ast.ImportFrom) and node.module:
            top_level_imports.append(node.module.split(".")[0])

    disallowed = {"anthropic", "openai", "pypdfium2"}
    assert not (disallowed & set(top_level_imports))


def test_ac9_validator_no_top_level_sdk_or_pypdfium2_import() -> None:
    """INV-5: validator.py 최상단에 anthropic/openai/pypdfium2 import가 없다(AST)."""
    src_path = Path(__file__).parent.parent / "src" / "pptx_md" / "validator.py"
    tree = ast.parse(src_path.read_text(encoding="utf-8"))

    top_level_imports: list[str] = []
    for node in tree.body:
        if isinstance(node, ast.Import):
            for alias in node.names:
                top_level_imports.append(alias.name.split(".")[0])
        elif isinstance(node, ast.ImportFrom) and node.module:
            top_level_imports.append(node.module.split(".")[0])

    disallowed = {"anthropic", "openai", "pypdfium2"}
    assert not (disallowed & set(top_level_imports))


def test_ac9_core_only_import_and_convert_succeed(tmp_path: Path) -> None:
    """INV-5: core-only 환경 프록시 -- import pptx_md + convert()가 정상 동작한다."""
    from pptx import Presentation as PptxPresentation

    from pptx_md.api import convert

    pptx_file = tmp_path / "deck.pptx"
    prs = PptxPresentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Core Only"
    prs.save(str(pptx_file))

    result = convert(pptx_file)
    assert isinstance(result, str)
    assert "Core Only" in result
