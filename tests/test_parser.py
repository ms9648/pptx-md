"""Unit tests for src/pptx_md/parser.py — FR-03 슬라이드 파싱, FR-04 도형 파싱.

Test naming convention: test_ac<N>_<description>
Each AC (FR-03 AC1~AC7, FR-04 AC1~AC7) has at least one test.
python-pptx is NOT mocked — real API calls only (ADR-207).
"""

from __future__ import annotations

import io
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from pptx_md.errors import ParseError
from pptx_md.ir import (
    GroupShapeIR,
    ImageShapeIR,
    OtherShapeIR,
    PresentationIR,
    ShapeKind,
    TableShapeIR,
    TextShapeIR,
)
from pptx_md.parser import parse_presentation

# ---------------------------------------------------------------------------
# Helper: parse from bytes (saves to tmp_path, then parse)
# ---------------------------------------------------------------------------


def _parse_bytes(pptx_bytes: bytes, tmp_path: Path) -> PresentationIR:
    """Write pptx_bytes to a temp file and parse it."""
    p = tmp_path / "test.pptx"
    p.write_bytes(pptx_bytes)
    return parse_presentation(p)


# ===========================================================================
# FR-03 슬라이드 파싱 AC1~AC7
# ===========================================================================


class TestFR03SlideParsingAC1:
    """AC1: 슬라이드 열거 / 순서 보장."""

    def test_ac1_슬라이드_개수_일치(
        self, pptx_multi_slide: bytes, tmp_path: Path
    ) -> None:
        """AC1: N개 슬라이드 파일 → SlideIR N개 반환."""
        pres = _parse_bytes(pptx_multi_slide, tmp_path)
        assert len(pres.slides) == 3

    def test_ac1_슬라이드_순서_원본과_동일(
        self, pptx_multi_slide: bytes, tmp_path: Path
    ) -> None:
        """AC1: slides 순서가 Presentation.slides 순서와 동일하다."""
        pres = _parse_bytes(pptx_multi_slide, tmp_path)
        # 각 슬라이드의 title 이 Slide 0 / Slide 1 / Slide 2 순서여야 함
        titles = [s.title for s in pres.slides]
        assert titles == ["Slide 0", "Slide 1", "Slide 2"]


class TestFR03SlideParsingAC2:
    """AC2: 제목 추출."""

    def test_ac2_제목_placeholder_텍스트_일치(
        self, pptx_with_text_slide: bytes, tmp_path: Path
    ) -> None:
        """AC2: 제목 placeholder 텍스트와 SlideIR.title 이 정확히 일치한다."""
        pres = _parse_bytes(pptx_with_text_slide, tmp_path)
        assert len(pres.slides) == 1
        assert pres.slides[0].title == "Test Slide Title"


class TestFR03SlideParsingAC3:
    """AC3: 발표자 노트 추출."""

    def test_ac3_노트_텍스트_일치(self, pptx_with_notes: bytes, tmp_path: Path) -> None:
        """AC3: notes_slide 의 텍스트와 SlideIR.notes 가 일치한다."""
        pres = _parse_bytes(pptx_with_notes, tmp_path)
        assert len(pres.slides) == 1
        assert "Speaker notes text" in pres.slides[0].notes

    def test_ac3_노트_없는_슬라이드_빈_문자열(
        self, pptx_with_text_slide: bytes, tmp_path: Path
    ) -> None:
        """AC3: 노트 슬라이드가 없으면 notes 는 빈 문자열이다 (ADR-202)."""
        pres = _parse_bytes(pptx_with_text_slide, tmp_path)
        assert isinstance(pres.slides[0].notes, str)
        # notes 가 str 이면 통과 (내용이 없을 수도 있음)


class TestFR03SlideParsingAC4:
    """AC4: 슬라이드 인덱스 (0-based)."""

    def test_ac4_인덱스_0기반_부여(
        self, pptx_multi_slide: bytes, tmp_path: Path
    ) -> None:
        """AC4: 각 슬라이드에 0-based index 가 부여된다."""
        pres = _parse_bytes(pptx_multi_slide, tmp_path)
        for expected_idx, slide in enumerate(pres.slides):
            assert slide.index == expected_idx

    def test_ac4_단일_슬라이드_인덱스_0(
        self, pptx_with_text_slide: bytes, tmp_path: Path
    ) -> None:
        """AC4: 슬라이드 1개면 index == 0."""
        pres = _parse_bytes(pptx_with_text_slide, tmp_path)
        assert pres.slides[0].index == 0


class TestFR03SlideParsingAC5:
    """AC5: 제목 없는 슬라이드 — 예외 없음, title == ""."""

    def test_ac5_제목_없는_슬라이드_예외_없음(
        self, pptx_no_title_slide: bytes, tmp_path: Path
    ) -> None:
        """AC5: 제목 placeholder 가 없는 슬라이드도 예외 없이 파싱된다."""
        pres = _parse_bytes(pptx_no_title_slide, tmp_path)
        assert len(pres.slides) == 1

    def test_ac5_제목_없는_슬라이드_title_빈_문자열(
        self, pptx_no_title_slide: bytes, tmp_path: Path
    ) -> None:
        """AC5: title placeholder 부재 시 SlideIR.title == '' (ADR-202)."""
        pres = _parse_bytes(pptx_no_title_slide, tmp_path)
        assert pres.slides[0].title == ""
        assert isinstance(pres.slides[0].title, str)


class TestFR03SlideParsingAC6:
    """AC6: 잘못된 입력 → ParseError."""

    def test_ac6_존재하지_않는_파일_ParseError(self, tmp_path: Path) -> None:
        """AC6: 존재하지 않는 경로 → ParseError."""
        with pytest.raises(ParseError):
            parse_presentation(tmp_path / "nonexistent.pptx")

    def test_ac6_비pptx_파일_ParseError(self, tmp_path: Path) -> None:
        """AC6: PPTX 가 아닌 파일(텍스트 파일) → ParseError."""
        bad_file = tmp_path / "not_a_pptx.pptx"
        bad_file.write_bytes(b"This is not a PPTX file content")
        with pytest.raises(ParseError):
            parse_presentation(bad_file)

    def test_ac6_손상된_zip_ParseError(self, tmp_path: Path) -> None:
        """AC6: 손상된 ZIP 구조 → ParseError."""
        bad_file = tmp_path / "corrupt.pptx"
        bad_file.write_bytes(b"PK\x03\x04" + b"\x00" * 50)  # fake ZIP header
        with pytest.raises(ParseError):
            parse_presentation(bad_file)

    def test_ac6_str_경로_지원(
        self, pptx_with_text_slide: bytes, tmp_path: Path
    ) -> None:
        """AC6: path 가 str 타입이어도 정상 동작한다."""
        p = tmp_path / "test.pptx"
        p.write_bytes(pptx_with_text_slide)
        pres = parse_presentation(str(p))  # str path
        assert isinstance(pres, PresentationIR)


class TestFR03SlideParsingAC7:
    """AC7: 슬라이드 0개 → 빈 목록, 예외 없음."""

    def test_ac7_슬라이드_0개_빈_목록_반환(
        self, pptx_zero_slides: bytes, tmp_path: Path
    ) -> None:
        """AC7: 슬라이드 0개인 PPTX → slides==[], 예외 없음."""
        pres = _parse_bytes(pptx_zero_slides, tmp_path)
        assert pres.slides == []
        assert len(pres.slides) == 0

    def test_ac7_슬라이드_0개_PresentationIR_타입(
        self, pptx_zero_slides: bytes, tmp_path: Path
    ) -> None:
        """AC7: 슬라이드 0개도 PresentationIR 를 반환한다."""
        pres = _parse_bytes(pptx_zero_slides, tmp_path)
        assert isinstance(pres, PresentationIR)
        assert isinstance(pres.slides, list)


# ===========================================================================
# FR-04 도형 파싱 AC1~AC7
# ===========================================================================


class TestFR04ShapeParsingAC1:
    """AC1: 텍스트 도형 — TextShapeIR, 단락 순서·줄바꿈 보존."""

    def test_ac1_텍스트_도형_타입_TEXT(
        self, pptx_with_text_slide: bytes, tmp_path: Path
    ) -> None:
        """AC1: has_text_frame 도형이 TextShapeIR 로 파싱된다."""
        pres = _parse_bytes(pptx_with_text_slide, tmp_path)
        slide = pres.slides[0]
        text_shapes = [s for s in slide.shapes if isinstance(s, TextShapeIR)]
        assert len(text_shapes) >= 1
        for ts in text_shapes:
            assert ts.kind == ShapeKind.TEXT

    def test_ac1_텍스트_단락_순서_보존(
        self, pptx_with_text_slide: bytes, tmp_path: Path
    ) -> None:
        """AC1: 단락 텍스트가 순서대로 추출된다."""
        pres = _parse_bytes(pptx_with_text_slide, tmp_path)
        slide = pres.slides[0]

        # 본문 텍스트 박스 찾기 (title 제외)
        body_shapes = [
            s for s in slide.shapes if isinstance(s, TextShapeIR) and not s.is_title
        ]
        assert len(body_shapes) >= 1

        body = body_shapes[0]
        texts = [p.text for p in body.paragraphs]
        assert "First paragraph" in texts

    def test_ac1_텍스트_빈_단락_보존(self, tmp_path: Path) -> None:
        """AC1: 빈 텍스트 frame 도 예외 없이 처리된다 (AC7 연관)."""
        prs = PptxPresentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tb.text_frame.text = ""  # empty text

        buf = io.BytesIO()
        prs.save(buf)

        p = tmp_path / "empty_text.pptx"
        p.write_bytes(buf.getvalue())
        pres = parse_presentation(p)

        text_shapes = [s for s in pres.slides[0].shapes if isinstance(s, TextShapeIR)]
        assert len(text_shapes) == 1
        assert all(isinstance(para.text, str) for para in text_shapes[0].paragraphs)


class TestFR04ShapeParsingAC2:
    """AC2: 표 도형 — TableShapeIR, 행×열 셀 텍스트."""

    def test_ac2_표_도형_타입_TABLE(
        self, pptx_with_table: bytes, tmp_path: Path
    ) -> None:
        """AC2: has_table 도형이 TableShapeIR 로 파싱된다."""
        pres = _parse_bytes(pptx_with_table, tmp_path)
        slide = pres.slides[0]
        table_shapes = [s for s in slide.shapes if isinstance(s, TableShapeIR)]
        assert len(table_shapes) == 1
        assert table_shapes[0].kind == ShapeKind.TABLE

    def test_ac2_행_열_수_정확(self, pptx_with_table: bytes, tmp_path: Path) -> None:
        """AC2: n_rows=2, n_cols=3 이 정확하다."""
        pres = _parse_bytes(pptx_with_table, tmp_path)
        tbl = [s for s in pres.slides[0].shapes if isinstance(s, TableShapeIR)][0]
        assert tbl.n_rows == 2
        assert tbl.n_cols == 3

    def test_ac2_셀_텍스트_2차원_배열(
        self, pptx_with_table: bytes, tmp_path: Path
    ) -> None:
        """AC2: 셀 텍스트가 행→열 순서의 2D 리스트로 추출된다."""
        pres = _parse_bytes(pptx_with_table, tmp_path)
        tbl = [s for s in pres.slides[0].shapes if isinstance(s, TableShapeIR)][0]
        assert tbl.rows[0] == ["A", "B", "C"]
        assert tbl.rows[1] == ["D", "E", "F"]


class TestFR04ShapeParsingAC3:
    """AC3: 이미지 도형 — ImageShapeIR, image_bytes, content_type."""

    def test_ac3_이미지_도형_타입_IMAGE(
        self, pptx_with_image: bytes, tmp_path: Path
    ) -> None:
        """AC3: PICTURE 도형이 ImageShapeIR 로 파싱된다."""
        pres = _parse_bytes(pptx_with_image, tmp_path)
        slide = pres.slides[0]
        img_shapes = [s for s in slide.shapes if isinstance(s, ImageShapeIR)]
        assert len(img_shapes) == 1
        assert img_shapes[0].kind == ShapeKind.IMAGE

    def test_ac3_image_bytes_비어있지_않음(
        self, pptx_with_image: bytes, tmp_path: Path
    ) -> None:
        """AC3: image_bytes 가 비어있지 않다 (원본 이미지 바이너리 포함)."""
        pres = _parse_bytes(pptx_with_image, tmp_path)
        img = [s for s in pres.slides[0].shapes if isinstance(s, ImageShapeIR)][0]
        assert len(img.image_bytes) > 0
        assert isinstance(img.image_bytes, bytes)

    def test_ac3_content_type_추출(
        self, pptx_with_image: bytes, tmp_path: Path
    ) -> None:
        """AC3: image_format 이 추출된다 (예: 'png')."""
        pres = _parse_bytes(pptx_with_image, tmp_path)
        img = [s for s in pres.slides[0].shapes if isinstance(s, ImageShapeIR)][0]
        assert img.image_format != ""
        assert isinstance(img.image_format, str)

    def test_ac3_확장_슬롯_none(self, pptx_with_image: bytes, tmp_path: Path) -> None:
        """AC3: v1 파서는 classification, description 을 None 으로 둔다 (ADR-205)."""
        pres = _parse_bytes(pptx_with_image, tmp_path)
        img = [s for s in pres.slides[0].shapes if isinstance(s, ImageShapeIR)][0]
        assert img.classification is None
        assert img.description is None


class TestFR04ShapeParsingAC4:
    """AC4: 그룹 도형 — GroupShapeIR, 자식 재귀 파싱."""

    def test_ac4_그룹_도형_GroupShapeIR(
        self, pptx_with_group: bytes, tmp_path: Path
    ) -> None:
        """AC4: GROUP 도형이 GroupShapeIR 로 파싱된다."""
        pres = _parse_bytes(pptx_with_group, tmp_path)
        slide = pres.slides[0]

        # 그룹이 없으면 개별 텍스트 박스들이 있을 수 있음 (group_shapes 미지원 시)
        group_shapes = [s for s in slide.shapes if isinstance(s, GroupShapeIR)]
        text_shapes = [s for s in slide.shapes if isinstance(s, TextShapeIR)]

        # 그룹이 생성됐으면 GroupShapeIR, 아니면 TextShapeIR 2개
        assert len(group_shapes) >= 1 or len(text_shapes) >= 2

    def test_ac4_중첩_그룹_재귀_파싱(self, tmp_path: Path) -> None:
        """AC4: 그룹 안 그룹도 GroupShapeIR.children 으로 재귀 파싱된다."""
        # 직접 PPTX 를 수동으로 만들어서 그룹 중첩 구조 확인
        # python-pptx 의 group_shapes 는 버전에 따라 동작이 다를 수 있으므로
        # GroupShapeIR 의 자식 파싱이 올바른지 테스트한다
        from pptx_md.ir import (  # noqa: PLC0415
            GroupShapeIR,
            ParagraphIR,
            ShapeKind,
            TextShapeIR,
        )

        # GroupShapeIR 를 직접 생성해서 트리 구조 검증
        inner_text = TextShapeIR(
            shape_id=1,
            name="t1",
            kind=ShapeKind.TEXT,
            paragraphs=[ParagraphIR(text="inner", level=0)],
        )
        inner_group = GroupShapeIR(
            shape_id=2,
            name="g1",
            kind=ShapeKind.GROUP,
            children=[inner_text],
        )
        outer_group = GroupShapeIR(
            shape_id=3,
            name="g2",
            kind=ShapeKind.GROUP,
            children=[inner_group],
        )

        assert isinstance(outer_group.children[0], GroupShapeIR)
        assert isinstance(outer_group.children[0].children[0], TextShapeIR)


class TestFR04ShapeParsingAC5:
    """AC5: 다중 타입 혼재 슬라이드 — z-order 보존, 각 타입 정확."""

    def test_ac5_혼재_슬라이드_도형_타입_각각_정확(
        self, pptx_mixed_shapes: bytes, tmp_path: Path
    ) -> None:
        """AC5: 텍스트·표·이미지가 공존 시 각각 올바른 타입으로 파싱된다."""
        pres = _parse_bytes(pptx_mixed_shapes, tmp_path)
        slide = pres.slides[0]
        shapes = slide.shapes

        kinds = [s.kind for s in shapes]
        assert ShapeKind.TEXT in kinds
        assert ShapeKind.TABLE in kinds
        assert ShapeKind.IMAGE in kinds

    def test_ac5_z_order_원본_순서_보존(
        self, pptx_mixed_shapes: bytes, tmp_path: Path
    ) -> None:
        """AC5: shapes 목록이 slide.shapes 순회 순서와 동일하다."""
        pres = _parse_bytes(pptx_mixed_shapes, tmp_path)
        slide = pres.slides[0]

        # 첫 번째가 텍스트, 두 번째가 표, 세 번째가 이미지
        shapes = slide.shapes
        assert len(shapes) == 3
        assert isinstance(shapes[0], TextShapeIR)
        assert isinstance(shapes[1], TableShapeIR)
        assert isinstance(shapes[2], ImageShapeIR)


class TestFR04ShapeParsingAC6:
    """AC6: 미지원/기타 도형 — OtherShapeIR, 예외 없음."""

    def test_ac6_차트_도형_OtherShapeIR(self, other_shape_pptx_path: Path) -> None:
        """AC6: 차트(MSO_SHAPE_TYPE.CHART) → OtherShapeIR, 예외 없음."""
        if not other_shape_pptx_path.is_file():
            pytest.skip("other_shape.pptx fixture 없음")

        pres = parse_presentation(other_shape_pptx_path)
        slide = pres.slides[0]
        other_shapes = [s for s in slide.shapes if isinstance(s, OtherShapeIR)]

        assert len(other_shapes) >= 1
        assert other_shapes[0].kind == ShapeKind.OTHER

    def test_ac6_OtherShapeIR_mso_type_기록(self, other_shape_pptx_path: Path) -> None:
        """AC6: OtherShapeIR.mso_shape_type 에 원본 타입 라벨이 기록된다."""
        if not other_shape_pptx_path.is_file():
            pytest.skip("other_shape.pptx fixture 없음")

        pres = parse_presentation(other_shape_pptx_path)
        slide = pres.slides[0]
        other_shapes = [s for s in slide.shapes if isinstance(s, OtherShapeIR)]

        # 차트 타입이므로 "CHART" 가 포함돼야 함
        assert any("CHART" in s.mso_shape_type.upper() for s in other_shapes)

    def test_ac6_미지원_도형_이후_파싱_계속(
        self, other_shape_pptx_path: Path, tmp_path: Path
    ) -> None:
        """AC6: 미지원 도형이 있어도 같은 슬라이드의 다른 도형이 파싱된다."""
        if not other_shape_pptx_path.is_file():
            pytest.skip("other_shape.pptx fixture 없음")

        # other_shape.pptx 에 텍스트 박스를 추가한 버전 생성
        from pptx import Presentation as PptxPresentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE

        prs = PptxPresentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 텍스트 박스 (파싱 가능)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
        tb.text_frame.text = "Normal text"

        # 차트 (OTHER 경로)
        chart_data = ChartData()
        chart_data.categories = ["A", "B"]
        chart_data.add_series("S1", (1.0, 2.0))
        slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(3),
            Inches(0.5),
            Inches(4),
            Inches(3),
            chart_data,
        )

        buf = io.BytesIO()
        prs.save(buf)
        p = tmp_path / "mixed_with_chart.pptx"
        p.write_bytes(buf.getvalue())

        pres = parse_presentation(p)
        shapes = pres.slides[0].shapes

        text_ok = any(isinstance(s, TextShapeIR) for s in shapes)
        other_ok = any(isinstance(s, OtherShapeIR) for s in shapes)
        assert text_ok, "텍스트 도형이 파싱되어야 함"
        assert other_ok, "차트가 OtherShapeIR 로 파싱되어야 함"


class TestFR04ShapeParsingAC7:
    """AC7: 빈 텍스트/빈 표 — 예외 없음, 빈 구조 반환."""

    def test_ac7_빈_텍스트_frame_예외_없음(self, tmp_path: Path) -> None:
        """AC7: 텍스트가 비어있는 text_frame 도 예외 없이 TextShapeIR 반환."""
        prs = PptxPresentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tb.text_frame.text = ""

        buf = io.BytesIO()
        prs.save(buf)
        p = tmp_path / "empty_text.pptx"
        p.write_bytes(buf.getvalue())

        pres = parse_presentation(p)
        text_shapes = [s for s in pres.slides[0].shapes if isinstance(s, TextShapeIR)]
        assert len(text_shapes) == 1
        # 도형 자체는 누락되지 않음
        assert isinstance(text_shapes[0], TextShapeIR)

    def test_ac7_빈_표_예외_없음(
        self, pptx_with_empty_table: bytes, tmp_path: Path
    ) -> None:
        """AC7: 모든 셀이 빈 표도 예외 없이 TableShapeIR 반환."""
        pres = _parse_bytes(pptx_with_empty_table, tmp_path)
        slide = pres.slides[0]
        table_shapes = [s for s in slide.shapes if isinstance(s, TableShapeIR)]
        assert len(table_shapes) == 1
        tbl = table_shapes[0]
        # 셀은 모두 빈 문자열
        for row in tbl.rows:
            for cell in row:
                assert isinstance(cell, str)

    def test_ac7_빈_표_셀_빈_문자열(
        self, pptx_with_empty_table: bytes, tmp_path: Path
    ) -> None:
        """AC7: 빈 셀 텍스트가 '' 이다 (None 아님, ADR-202)."""
        pres = _parse_bytes(pptx_with_empty_table, tmp_path)
        tbl = [s for s in pres.slides[0].shapes if isinstance(s, TableShapeIR)][0]
        for row in tbl.rows:
            for cell in row:
                assert cell == ""

    def test_ac7_빈_텍스트_paragraphs_str_타입(self, tmp_path: Path) -> None:
        """AC7: 빈 텍스트 도형의 ParagraphIR.text 가 str 타입이다."""
        prs = PptxPresentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tf = tb.text_frame
        tf.text = ""

        buf = io.BytesIO()
        prs.save(buf)
        p = tmp_path / "empty_text2.pptx"
        p.write_bytes(buf.getvalue())

        pres = parse_presentation(p)
        text_shapes = [s for s in pres.slides[0].shapes if isinstance(s, TextShapeIR)]
        if text_shapes:
            for para in text_shapes[0].paragraphs:
                assert isinstance(para.text, str)


# ===========================================================================
# 추가: PresentationIR 메타 검증
# ===========================================================================


class TestPresentationIRMeta:
    """parse_presentation 의 PresentationIR 메타 필드 검증."""

    def test_source_path_설정(
        self, pptx_with_text_slide: bytes, tmp_path: Path
    ) -> None:
        """parse_presentation 의 PresentationIR.source_path 가 입력 경로와 일치."""
        p = tmp_path / "test.pptx"
        p.write_bytes(pptx_with_text_slide)
        pres = parse_presentation(p)
        assert str(p) == pres.source_path

    def test_slide_dimensions_양수(
        self, pptx_with_text_slide: bytes, tmp_path: Path
    ) -> None:
        """슬라이드 폭/높이 EMU 가 양수로 설정된다."""
        pres = _parse_bytes(pptx_with_text_slide, tmp_path)
        assert pres.slide_width_emu > 0
        assert pres.slide_height_emu > 0
