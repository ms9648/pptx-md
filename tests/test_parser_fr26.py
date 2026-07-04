"""Unit tests for FR-26 (#67): Parser/COM 커버리지 확장.

Test naming convention: test_ac<N>_<description> (skill §5).
python-pptx is NOT mocked for shape trees — real API calls only (ADR-207).
Mocks are used only where python-pptx offers no way to reproduce an edge
case (e.g. paragraph.level == None, per-shape exceptions).

Scope confirmed with PM prior to implementation (see issue #67):
- AC1/AC2 (Group recursion) — already implemented in _parse_group; tests
  here are regression coverage focused on iter_shapes() flattening.
- AC3/AC4 (paragraph.level capture) — already implemented in _parse_text;
  tests here focus on accuracy across mixed levels and None defense.
- AC5/AC6 (Chart/SmartArt fallback_text / mso_shape_type) — the actual
  implementation gap closed by this issue, in _parse_other() /
  _extract_graphic_frame_fallback_text() / _graphic_frame_type_label().
- AC7 (partial-failure isolation) and AC10 (determinism) — regression tests.
- AC8 (no image loss) — golden-fixture regression test.
- AC9 (COM path) — N/A, see TestFR26COMPathNotApplicableAC9 below.
"""

from __future__ import annotations

import io
import types
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from pptx_md.ir import (
    GroupShapeIR,
    ImageShapeIR,
    OtherShapeIR,
    PresentationIR,
    ShapeKind,
    TableShapeIR,
    TextShapeIR,
    iter_shapes,
)
from pptx_md.parser import (
    _extract_graphic_frame_fallback_text,
    _graphic_frame_type_label,
    _parse_slide,
    _parse_text,
    parse_presentation,
)

# ---------------------------------------------------------------------------
# Helper
# ---------------------------------------------------------------------------


def _parse_bytes(
    pptx_bytes: bytes, tmp_path: Path, filename: str = "t.pptx"
) -> PresentationIR:
    p = tmp_path / filename
    p.write_bytes(pptx_bytes)
    return parse_presentation(p)


# ===========================================================================
# AC1: Group(자식3: 텍스트·표·그림) — children 길이 3, 타입 정확, iter_shapes 평탄화
# ===========================================================================


class TestFR26AC1GroupThreeChildren:
    def test_ac1_그룹_자식_3개_타입_정확(
        self, pptx_group_three_children: bytes, tmp_path: Path
    ) -> None:
        """ac1_그룹_자식_3개_타입_정확: children 길이 3, 각 자식 올바른 IR 타입."""
        pres = _parse_bytes(pptx_group_three_children, tmp_path)
        groups = [s for s in pres.slides[0].shapes if isinstance(s, GroupShapeIR)]
        assert len(groups) == 1
        grp = groups[0]
        assert len(grp.children) == 3
        assert isinstance(grp.children[0], TextShapeIR)
        assert isinstance(grp.children[1], TableShapeIR)
        assert isinstance(grp.children[2], ImageShapeIR)

    def test_ac1_iter_shapes_자식_3개_방출(
        self, pptx_group_three_children: bytes, tmp_path: Path
    ) -> None:
        """ac1_iter_shapes_자식_3개_방출: 평탄화 시 group + 자식 3개 = 4개."""
        pres = _parse_bytes(pptx_group_three_children, tmp_path)
        slide = pres.slides[0]
        flattened = list(iter_shapes(slide))
        # group 노드 자신(1) + children(3) = 4
        assert len(flattened) == 4
        leaf_kinds = [s.kind for s in flattened if not isinstance(s, GroupShapeIR)]
        assert leaf_kinds.count(ShapeKind.TEXT) == 1
        assert leaf_kinds.count(ShapeKind.TABLE) == 1
        assert leaf_kinds.count(ShapeKind.IMAGE) == 1


# ===========================================================================
# AC2: 중첩 Group(2단) — 리프까지 재귀, iter_shapes 방출 수 == 트리 전체 노드 수
# ===========================================================================


class TestFR26AC2NestedGroupTwoLevels:
    def test_ac2_중첩_그룹_리프까지_재귀(
        self, pptx_nested_group_two_levels: bytes, tmp_path: Path
    ) -> None:
        """ac2_중첩_그룹_리프까지_재귀: outer>inner>leaf 구조가 예외 없이 파싱된다."""
        pres = _parse_bytes(pptx_nested_group_two_levels, tmp_path)
        slide = pres.slides[0]

        outer_groups = [s for s in slide.shapes if isinstance(s, GroupShapeIR)]
        assert len(outer_groups) == 1
        outer = outer_groups[0]

        inner_groups = [c for c in outer.children if isinstance(c, GroupShapeIR)]
        assert len(inner_groups) == 1
        inner = inner_groups[0]

        leaf_texts = [c for c in inner.children if isinstance(c, TextShapeIR)]
        assert len(leaf_texts) == 1
        assert "deep leaf text" in leaf_texts[0].paragraphs[0].text

    def test_ac2_iter_shapes_방출수_트리전체노드수_일치(
        self, pptx_nested_group_two_levels: bytes, tmp_path: Path
    ) -> None:
        """ac2_iter_shapes_방출수_일치: outer/inner/leaf/sibling 4개 노드가 방출된다."""
        pres = _parse_bytes(pptx_nested_group_two_levels, tmp_path)
        slide = pres.slides[0]

        flattened = list(iter_shapes(slide))
        # outer group(1) + inner group(1) + leaf text(1) + sibling text(1) = 4
        assert len(flattened) == 4
        assert sum(1 for s in flattened if isinstance(s, GroupShapeIR)) == 2
        assert sum(1 for s in flattened if isinstance(s, TextShapeIR)) == 2


# ===========================================================================
# AC3: 0/1/2 레벨 혼합 불릿 — ParagraphIR.level 이 원본과 일치
# ===========================================================================


class TestFR26AC3ParagraphLevelAccuracy:
    def test_ac3_불릿_레벨_0_1_2_정확히_일치(self, tmp_path: Path) -> None:
        """ac3_불릿_레벨_정확히_일치: 0/1/2 레벨 단락이 각각 원본 정수와 일치한다."""
        prs = PptxPresentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
        tf = tb.text_frame
        tf.text = "level 0"
        tf.paragraphs[0].level = 0
        p1 = tf.add_paragraph()
        p1.text = "level 1"
        p1.level = 1
        p2 = tf.add_paragraph()
        p2.text = "level 2"
        p2.level = 2

        buf = io.BytesIO()
        prs.save(buf)
        p = tmp_path / "levels.pptx"
        p.write_bytes(buf.getvalue())

        pres = parse_presentation(p)
        text_shapes = [s for s in pres.slides[0].shapes if isinstance(s, TextShapeIR)]
        assert len(text_shapes) == 1
        paragraphs = text_shapes[0].paragraphs
        assert [pg.level for pg in paragraphs] == [0, 1, 2]
        assert [pg.text for pg in paragraphs] == ["level 0", "level 1", "level 2"]


# ===========================================================================
# AC4: paragraph.level 이 None → ParagraphIR.level == 0 정규화, 예외 없음
# ===========================================================================


class TestFR26AC4ParagraphLevelNoneDefense:
    def test_ac4_level_none_이면_0으로_정규화(self) -> None:
        """ac4_level_none_0정규화: level이 None이어도 예외 없이 0으로 정규화된다."""
        from unittest.mock import MagicMock

        mock_para = MagicMock()
        mock_para.level = None
        mock_para.runs = []
        mock_para.text = "no level info"

        mock_shape = MagicMock()
        mock_shape.placeholder_format = None
        mock_shape.text_frame.paragraphs = [mock_para]

        result = _parse_text(mock_shape, sid=1, name="none_level_test")
        assert len(result.paragraphs) == 1
        assert result.paragraphs[0].level == 0
        assert result.paragraphs[0].text == "no level info"


# ===========================================================================
# AC5/AC6: Chart/SmartArt — fallback_text 비공백 보존, mso_shape_type 실제 라벨
# ===========================================================================


class TestFR26AC5AC6ChartSmartArtFallback:
    def test_ac5_차트_제목_fallback_text_보존(
        self, pptx_chart_with_title: bytes, tmp_path: Path
    ) -> None:
        """ac5_차트_제목_보존: 차트 제목이 fallback_text 에 비공백으로 보존된다."""
        pres = _parse_bytes(pptx_chart_with_title, tmp_path)
        others = [s for s in pres.slides[0].shapes if isinstance(s, OtherShapeIR)]
        assert len(others) == 1
        assert others[0].fallback_text != ""
        assert "Sales Chart" in others[0].fallback_text

    def test_ac6_차트_mso_shape_type_CHART(
        self, pptx_chart_with_title: bytes, tmp_path: Path
    ) -> None:
        """ac6_차트_mso_type_CHART: 차트의 mso_shape_type 이 'CHART' (≠UNKNOWN)."""
        pres = _parse_bytes(pptx_chart_with_title, tmp_path)
        others = [s for s in pres.slides[0].shapes if isinstance(s, OtherShapeIR)]
        assert len(others) == 1
        assert others[0].mso_shape_type == "CHART"
        assert others[0].mso_shape_type != "UNKNOWN"

    def test_ac5_차트_카테고리_시리즈_캐시_텍스트_보존(
        self, other_shape_pptx_path: Path
    ) -> None:
        """ac5_카테고리_시리즈_보존: 제목 없이도 캐시 문자열이 보존된다."""
        if not other_shape_pptx_path.is_file():
            pytest.skip("other_shape.pptx fixture 없음")
        pres = parse_presentation(other_shape_pptx_path)
        others = [s for s in pres.slides[0].shapes if isinstance(s, OtherShapeIR)]
        assert len(others) >= 1
        chart_other = others[0]
        assert chart_other.fallback_text != ""
        # tests/fixtures/other_shape.pptx 는 East/West/Midwest 카테고리를 담고 있다
        assert "East" in chart_other.fallback_text

    def test_ac5_smartart_스텁_노드_텍스트_보존(
        self, pptx_smartart_stub: bytes, tmp_path: Path
    ) -> None:
        """ac5_smartart_노드텍스트_보존: 인라인 a:t 텍스트가 보존된다."""
        pres = _parse_bytes(pptx_smartart_stub, tmp_path)
        others = [s for s in pres.slides[0].shapes if isinstance(s, OtherShapeIR)]
        assert len(others) == 1
        assert "Node Alpha" in others[0].fallback_text
        assert "Node Beta" in others[0].fallback_text

    def test_ac6_smartart_mso_shape_type_DIAGRAM(
        self, pptx_smartart_stub: bytes, tmp_path: Path
    ) -> None:
        """ac6_smartart_mso_type_DIAGRAM: mso_shape_type=='DIAGRAM' (≠UNKNOWN)."""
        pres = _parse_bytes(pptx_smartart_stub, tmp_path)
        others = [s for s in pres.slides[0].shapes if isinstance(s, OtherShapeIR)]
        assert len(others) == 1
        assert others[0].mso_shape_type == "DIAGRAM"
        assert others[0].mso_shape_type != "UNKNOWN"

    def test_ac5_추출_불가_시_빈문자열_무예외(self) -> None:
        """ac5_추출불가_빈문자열_무예외: chart 접근 실패해도 예외 없이 ''."""

        class _BrokenChartShape:
            has_chart = True
            element = None

            @property
            def chart(self) -> object:
                raise RuntimeError("chart part unavailable")

        # 예외가 전파되지 않아야 한다 (ADR-204)
        result = _extract_graphic_frame_fallback_text(_BrokenChartShape())
        assert result == ""

    def test_ac6_graphicData_없는_shape_None_라벨(self) -> None:
        """ac6_graphicData_없음_None: graphicData 를 찾을 수 없으면 None 을 반환한다."""

        class _NoGraphicDataShape:
            element = None

        assert _graphic_frame_type_label(_NoGraphicDataShape()) is None

    def test_ac6_graphicData_노드_부재시_None(self) -> None:
        """ac6_graphicData_노드_부재: graphicData 자식이 없으면 None."""
        from lxml import etree  # noqa: PLC0415

        class _ShapeNoGraphicDataChild:
            element = etree.Element("root")

        assert _graphic_frame_type_label(_ShapeNoGraphicDataChild()) is None

    def test_ac6_find_예외시_None_무예외(self) -> None:
        """ac6_find_예외_None: element.find 가 예외를 던져도 무예외로 None 반환."""

        class _RaisingFindElement:
            def find(self, path: str) -> object:
                raise RuntimeError("boom")

        class _ShapeWithRaisingFind:
            element = _RaisingFindElement()

        assert _graphic_frame_type_label(_ShapeWithRaisingFind()) is None

    def test_ac6_알수없는_graphicData_uri_GRAPHIC_FRAME_라벨(self) -> None:
        """ac6_미지_uri_GRAPHIC_FRAME: diagram 이 아닌 임의 uri 는 GRAPHIC_FRAME."""
        from lxml import etree  # noqa: PLC0415

        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        elem = etree.fromstring(
            f'<root xmlns:a="{a_ns}">'
            '<a:graphic><a:graphicData uri="http://example.com/unknown"/>'
            "</a:graphic></root>"
        )

        class _ShapeUnknownGraphicData:
            element = elem

        assert _graphic_frame_type_label(_ShapeUnknownGraphicData()) == "GRAPHIC_FRAME"

    def test_ac6_graphicData_uri_빈값_None(self) -> None:
        """ac6_uri_빈값_None: graphicData 는 있으나 uri 속성이 없으면 None."""
        from lxml import etree  # noqa: PLC0415

        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        elem = etree.fromstring(
            f'<root xmlns:a="{a_ns}"><a:graphic><a:graphicData/></a:graphic></root>'
        )

        class _ShapeEmptyUri:
            element = elem

        assert _graphic_frame_type_label(_ShapeEmptyUri()) is None

    def test_ac5_차트XML_추출중_예외_무예외_빈문자열(self) -> None:
        """ac5_차트XML_예외_무예외: chart.element.iter 예외 시 ''를 반환한다."""

        class _RaisingIterElement:
            def iter(self, tag: str) -> object:
                raise RuntimeError("malformed chart xml")

        class _ChartWithRaisingXml:
            has_chart = True
            element = None

            @property
            def chart(self) -> object:
                return types.SimpleNamespace(element=_RaisingIterElement())

        result = _extract_graphic_frame_fallback_text(_ChartWithRaisingXml())
        assert result == ""


# ===========================================================================
# AC7: 예외 던지는 손상 도형 1개가 정상 도형들 사이 — 해당만 강등, 좌표 유지
# ===========================================================================


class TestFR26AC7PartialFailureIsolation:
    def test_ac7_손상도형_주변_정상도형_유지_좌표보존(self, tmp_path: Path) -> None:
        """ac7_손상도형_격리: 손상 도형만 강등, 나머지 정상, 좌표 유지."""
        prs = PptxPresentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tb1.text_frame.text = "Normal A"
        tb2 = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(2), Inches(1))
        tb2.text_frame.text = "Normal B"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = PptxPresentation(buf)
        real_shapes = list(prs2.slides[0].shapes)

        class _CorruptedShape:
            shape_id = 999
            name = "Corrupted"
            shape_type = None
            has_table = False
            has_text_frame = True
            left = 111
            top = 222
            width = 333
            height = 444

            @property
            def placeholder_format(self) -> object:
                raise RuntimeError("boom: corrupted shape data")

        fake_slide = types.SimpleNamespace(
            shapes=[real_shapes[0], _CorruptedShape(), real_shapes[1]]
        )

        slide_ir = _parse_slide(0, fake_slide)
        assert len(slide_ir.shapes) == 3

        first, corrupted, last = slide_ir.shapes
        assert isinstance(first, TextShapeIR)
        assert isinstance(last, TextShapeIR)
        assert "Normal A" in "".join(p.text for p in first.paragraphs)
        assert "Normal B" in "".join(p.text for p in last.paragraphs)

        assert isinstance(corrupted, OtherShapeIR)
        assert corrupted.mso_shape_type == "UNKNOWN"
        # FR-23 좌표 필드는 강등 시에도 유지된다
        assert corrupted.left == 111
        assert corrupted.top == 222
        assert corrupted.width == 333
        assert corrupted.height == 444


# ===========================================================================
# AC8: 이미지 포함 골든 픽스처 — 생성 ImageShapeIR 수 기대값 일치, 그림 손실 0
# ===========================================================================


class TestFR26AC8NoImageLossAcrossGroups:
    def test_ac8_이미지_3개_손실_0(
        self, pptx_images_including_grouped: bytes, tmp_path: Path
    ) -> None:
        """ac8_이미지_손실_0: 그룹 내부 포함 3개 이미지가 모두 파싱된다."""
        pres = _parse_bytes(pptx_images_including_grouped, tmp_path)
        slide = pres.slides[0]

        images = [s for s in iter_shapes(slide) if isinstance(s, ImageShapeIR)]
        assert len(images) == 3
        for img in images:
            assert len(img.image_bytes) > 0


# ===========================================================================
# AC9: N/A — convert_via_com 미구현 (COM 경로 코드베이스에 부재)
# ===========================================================================


class TestFR26COMPathNotApplicableAC9:
    @pytest.mark.skip(
        reason=(
            "AC9 N/A: convert_via_com 은 코드베이스에 존재하지 않는다. "
            "이슈 #67 스코프에서 신규 COM 모듈 도입은 금지되었으며(발명 금지), "
            "COM 기반 변환 경로 구현은 별도 이슈로 이관되는 미구현 부채로 명시한다."
        )
    )
    def test_ac9_com_경로_미구현_부채(self) -> None:  # pragma: no cover
        """ac9_com_경로_미구현_부채: convert_via_com 부재 — 후속 이슈에서 다룰 부채."""
        raise AssertionError(
            "unreachable: this test documents an accepted debt via skip"
        )


# ===========================================================================
# AC10: 동일 입력 2회 파싱 — 도형 순서·필드값까지 동일 (결정성, ADR-218)
# ===========================================================================


class TestFR26AC10Determinism:
    def test_ac10_동일_입력_2회_파싱_완전동일(
        self, pptx_mixed_shapes: bytes, tmp_path: Path
    ) -> None:
        """ac10_결정성: 동일 PPTX 를 2회 파싱하면 전체 IR 트리가 완전히 동일하다."""
        p = tmp_path / "determinism.pptx"
        p.write_bytes(pptx_mixed_shapes)

        pres1 = parse_presentation(p)
        pres2 = parse_presentation(p)

        assert pres1 == pres2

    def test_ac10_동일_입력_2회_파싱_차트_스마트아트_포함_완전동일(
        self, pptx_chart_with_title: bytes, tmp_path: Path
    ) -> None:
        """ac10_결정성_차트포함: Chart/OtherShapeIR 강등 경로도 결정적으로 동일하다."""
        p = tmp_path / "determinism_chart.pptx"
        p.write_bytes(pptx_chart_with_title)

        pres1 = parse_presentation(p)
        pres2 = parse_presentation(p)

        assert pres1 == pres2
