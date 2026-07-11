"""Integration smoke test — real LibreOffice (soffice) -> PDF -> PNG -> 재구성
end-to-end (issue #95, ARCH-v020 §3.1/§3.4, 사람 승인 사항).

이 테스트는 실 VLM 은 호출하지 않는다(FakeReconstructor 주입 — 비용/네트워크
의존 없음). 검증 대상은 어디까지나 real ``slide_render.render_deck``
(soffice subprocess + pypdfium2 rasterisation)과 ``reconstruct_slides``
오케스트레이션의 배선이 실제로 맞물려 동작하는지이다.

환경 게이트(ARCH-v020 §3.7 AC6): ``slide_render_available()`` 이 False
(soffice 또는 pypdfium2 미설치)이면 이 테스트는 skip 된다. 이 개발 환경에는
soffice 가 없으므로 로컬에서는 skip 되고, soffice 가 설치된 CI 에서 실행된다.

픽스처: 합성 PPTX(python-pptx로 즉석 생성) 만 사용 — 사내 파일 0건(FR-35 AC7
정신 계승, RISK 반영).
"""

from __future__ import annotations

import io
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from pptx_md.slide_reconstruction import reconstruct_slides
from pptx_md.slide_render import slide_render_available

pytestmark = pytest.mark.skipif(
    not slide_render_available(),
    reason="soffice/pypdfium2 not available in this environment ([render] extra + "
    "LibreOffice install required) — skipped locally, runs in soffice-equipped CI",
)


class _FakeReconstructor:
    """No real VLM call — records the real PNG bytes it was handed."""

    def __init__(self) -> None:
        self.seen: list[tuple[int, int]] = []  # (slide_index, png byte length)

    def reconstruct(self, image_bytes: bytes, image_ext: str, context: object) -> str:
        idx = getattr(context, "slide_index", -1)
        self.seen.append((idx, len(image_bytes)))
        assert image_bytes[:8] == b"\x89PNG\r\n\x1a\n", "must be a real PNG render"
        assert image_ext == "png"
        return f"## Reconstructed slide {idx}\n\n(structured body)"


def _build_synthetic_multi_slide_pptx(path: Path) -> None:
    """A 2-slide synthetic deck: one table-bearing slide, one plain-text slide."""
    prs = PptxPresentation()

    table_slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbl = table_slide.shapes.add_table(
        2, 2, Inches(1), Inches(1), Inches(4), Inches(2)
    ).table
    tbl.cell(0, 0).text = "Col A"
    tbl.cell(0, 1).text = "Col B"
    tbl.cell(1, 0).text = "1"
    tbl.cell(1, 1).text = "2"

    text_slide = prs.slides.add_slide(prs.slide_layouts[1])
    if text_slide.shapes.title is not None:
        text_slide.shapes.title.text = "Plain Text Slide"

    buf = io.BytesIO()
    prs.save(buf)
    path.write_bytes(buf.getvalue())


def test_real_soffice_render_to_reconstruction_end_to_end(tmp_path: Path) -> None:
    """실 soffice(--convert-to pdf)+pypdfium2 렌더 -> reconstruct_slides 배선 통합
    검증."""
    from pptx_md.parser import parse_presentation

    pptx_path = tmp_path / "synthetic_deck.pptx"
    _build_synthetic_multi_slide_pptx(pptx_path)

    presentation = parse_presentation(pptx_path)
    assert len(presentation.slides) == 2

    reconstructor = _FakeReconstructor()
    # force_render=frozenset({0}) — deterministic regardless of the real
    # complexity heuristic's calibration; forces the *real* render_deck path
    # for slide 0 without depending on real VLM output.
    reconstruct_slides(
        presentation,
        reconstructor,
        force_render=frozenset({0}),
        force_text=frozenset({1}),
    )

    assert presentation.slides[0].reconstructed_md is not None
    assert "Reconstructed slide 0" in presentation.slides[0].reconstructed_md
    assert presentation.slides[1].reconstructed_md is None  # force_text
    assert reconstructor.seen, "reconstruct() must have been called with real PNG bytes"
    assert reconstructor.seen[0][1] > 0, "rendered PNG must be non-empty"
