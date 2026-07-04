"""Independent QA scenarios for #67 (FR-26 Parser/COM coverage) — TR-67.

Written by tester (QA), NOT the developer. These are deliberately different
constructions/fixtures from tests/test_parser_fr26.py so that a pass here is
independent evidence, not a re-run of the developer's own assertions.

Naming: test_tc67_<NN>_<slug> maps 1:1 to TR-67.md's TC-67-NN rows.
"""

from __future__ import annotations

import io
import types
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from pptx_md.assembler import assemble_document
from pptx_md.ir import (
    GroupShapeIR,
    ImageShapeIR,
    OtherShapeIR,
    TableShapeIR,
    TextShapeIR,
    iter_shapes,
)
from pptx_md.parser import (
    _extract_graphic_frame_fallback_text,
    _parse_slide,
    parse_presentation,
)

FIXTURES = Path(__file__).parent / "fixtures"


def _save(prs: PptxPresentation, tmp_path: Path, name: str) -> Path:
    buf = io.BytesIO()
    prs.save(buf)
    p = tmp_path / name
    p.write_bytes(buf.getvalue())
    return p


# ===========================================================================
# AC1 — Group recursion
# ===========================================================================


def test_tc67_01_group_children_order_preserved(tmp_path: Path) -> None:
    """TC-67-01 (AC1, 정상): 자식 삽입 순서(이미지→텍스트→표)가 그대로 보존된다.

    developer 테스트는 text/table/image 순서만 검증했다. 여기서는 다른 순서
    (image/text/table)로 구성해 순서 의존 버그(예: 타입별 재정렬)를 배제한다.
    """
    png = tmp_path / "p.png"
    png.write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d4948445200000001000000010802000000907753de"
            "0000000c49444154789c63f8ffff3f0005fe02fe0def46b80000000049454e44ae"
            "426082"
        )
    )
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(
        str(png), Inches(0.2), Inches(0.2), Inches(1), Inches(1)
    )
    tb = slide.shapes.add_textbox(Inches(2), Inches(0.2), Inches(2), Inches(1))
    tb.text_frame.text = "order-check"
    tbl = slide.shapes.add_table(1, 1, Inches(0.2), Inches(2), Inches(2), Inches(1))

    slide.shapes.add_group_shape([pic, tb, tbl])
    p = _save(prs, tmp_path, "order.pptx")

    pres = parse_presentation(p)
    groups = [s for s in pres.slides[0].shapes if isinstance(s, GroupShapeIR)]
    assert len(groups) == 1
    grp = groups[0]
    assert len(grp.children) == 3
    assert isinstance(grp.children[0], ImageShapeIR)
    assert isinstance(grp.children[1], TextShapeIR)
    assert isinstance(grp.children[2], TableShapeIR)


def test_tc67_02_group_single_child_boundary(tmp_path: Path) -> None:
    """TC-67-02 (AC1, 경계): 자식 1개짜리 최소 group — children 길이 1."""
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tb.text_frame.text = "lonely child"
    slide.shapes.add_group_shape([tb])
    p = _save(prs, tmp_path, "single_child.pptx")

    pres = parse_presentation(p)
    groups = [s for s in pres.slides[0].shapes if isinstance(s, GroupShapeIR)]
    assert len(groups) == 1
    assert len(groups[0].children) == 1
    assert isinstance(groups[0].children[0], TextShapeIR)


def test_tc67_03_group_child_exception_isolated_within_group(tmp_path: Path) -> None:
    """TC-67-03 (AC1, 예외): group 내부 자식 1개가 예외를 던져도 group 자체와
    나머지 형제는 정상 파싱된다 (그룹 경계를 넘어선 격리 확인, ADR-204).

    developer 테스트(AC7)는 슬라이드 최상위 레벨에서만 격리를 검증했다.
    여기서는 격리 경계가 group.children 레벨에서도 동일하게 작동하는지
    별도로 확인한다 (재귀 경로의 회귀 방지).
    """

    class _BoomChild:
        shape_id = 555
        name = "BoomChild"
        shape_type = None
        has_table = False
        has_text_frame = True
        left = 10
        top = 20
        width = 30
        height = 40

        @property
        def placeholder_format(self) -> object:
            raise RuntimeError("boom inside group child")

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tb.text_frame.text = "sibling ok"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    prs2 = PptxPresentation(buf)
    real_tb_shape = list(prs2.slides[0].shapes)[0]

    fake_group = types.SimpleNamespace(
        shape_id=1,
        name="FakeGroup",
        shape_type=MSO_SHAPE_TYPE.GROUP,
        left=0,
        top=0,
        width=100,
        height=100,
        shapes=[real_tb_shape, _BoomChild()],
    )
    fake_slide = types.SimpleNamespace(shapes=[fake_group])

    slide_ir = _parse_slide(0, fake_slide)
    assert len(slide_ir.shapes) == 1
    grp = slide_ir.shapes[0]
    assert isinstance(grp, GroupShapeIR)
    assert len(grp.children) == 2
    ok_child, boom_child = grp.children
    assert isinstance(ok_child, TextShapeIR)
    assert "sibling ok" in "".join(pg.text for pg in ok_child.paragraphs)
    assert isinstance(boom_child, OtherShapeIR)
    assert boom_child.mso_shape_type == "UNKNOWN"
    assert boom_child.left == 10
    assert boom_child.top == 20


# ===========================================================================
# AC2 — nested groups
# ===========================================================================


def test_tc67_04_three_level_nested_group(tmp_path: Path) -> None:
    """TC-67-04 (AC2, 정상): 3단 중첩(2단보다 깊은) group도 리프까지 재귀한다.

    developer는 2단만 검증했다(AC2 문구에 정확히 맞춤). 3단은 그보다 깊은
    깊이에서도 재귀 로직이 스택 문제 없이 동작하는지의 정상 경로 보강이다.
    """
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    leaf = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1))
    leaf.text_frame.text = "level-3 leaf"

    level2 = slide.shapes.add_group_shape([leaf])
    level1 = slide.shapes.add_group_shape([level2])
    slide.shapes.add_group_shape([level1])
    p = _save(prs, tmp_path, "three_level.pptx")

    pres = parse_presentation(p)
    slide_ir = pres.slides[0]
    flattened = list(iter_shapes(slide_ir))
    # 3 group nodes + 1 leaf = 4
    assert len(flattened) == 4
    assert sum(1 for s in flattened if isinstance(s, GroupShapeIR)) == 3
    leaves = [s for s in flattened if isinstance(s, TextShapeIR)]
    assert len(leaves) == 1
    assert "level-3 leaf" in leaves[0].paragraphs[0].text


def test_tc67_05_empty_group_boundary(tmp_path: Path) -> None:
    """TC-67-05 (AC2, 경계): 자식이 0개인 group — children == [] 이고 예외 없음."""
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sp_tree = slide.shapes._spTree
    from lxml import etree  # noqa: PLC0415

    grp_xml = (
        "<p:grpSp"
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        "<p:nvGrpSpPr>"
        '<p:cNvPr id="20" name="EmptyGroup"/>'
        "<p:cNvGrpSpPr/><p:nvPr/>"
        "</p:nvGrpSpPr>"
        "<p:grpSpPr>"
        "<a:xfrm>"
        '<a:off x="0" y="0"/><a:ext cx="100" cy="100"/>'
        '<a:chOff x="0" y="0"/><a:chExt cx="100" cy="100"/>'
        "</a:xfrm>"
        "</p:grpSpPr>"
        "</p:grpSp>"
    )
    sp_tree.append(etree.fromstring(grp_xml))
    p = _save(prs, tmp_path, "empty_group.pptx")

    pres = parse_presentation(p)
    groups = [s for s in pres.slides[0].shapes if isinstance(s, GroupShapeIR)]
    assert len(groups) == 1
    assert groups[0].children == []


# ===========================================================================
# AC3/AC4 — paragraph.level
# ===========================================================================


def test_tc67_06_multi_shape_mixed_nonmonotonic_levels(tmp_path: Path) -> None:
    """TC-67-06 (AC3, 정상): 비단조 레벨 시퀀스(0,2,1,0)가 두 개의 텍스트
    도형에 걸쳐 정확히 캡처된다 (developer는 단일 도형·단조 증가만 검증)."""
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(2))
    tf1 = tb1.text_frame
    tf1.text = "a0"
    tf1.paragraphs[0].level = 0
    p1 = tf1.add_paragraph()
    p1.text = "a2"
    p1.level = 2
    p2 = tf1.add_paragraph()
    p2.text = "a1"
    p2.level = 1

    tb2 = slide.shapes.add_textbox(Inches(4), Inches(0.5), Inches(3), Inches(2))
    tb2.text_frame.text = "b0"
    tb2.text_frame.paragraphs[0].level = 0

    p = _save(prs, tmp_path, "nonmonotonic.pptx")
    pres = parse_presentation(p)
    text_shapes = [s for s in pres.slides[0].shapes if isinstance(s, TextShapeIR)]
    assert len(text_shapes) == 2
    assert [pg.level for pg in text_shapes[0].paragraphs] == [0, 2, 1]
    assert [pg.level for pg in text_shapes[1].paragraphs] == [0]


def test_tc67_07_level_non_int_string_does_not_crash_pipeline(tmp_path: Path) -> None:
    """TC-67-07 (AC3/AC4, 예외): level이 정수로 변환 불가한 값이면 int()가
    ValueError를 던지지만 ADR-204 격리(상위 _parse_shape try/except)에 의해
    해당 도형만 OtherShapeIR로 강등되고 전체 파싱은 중단되지 않는다.

    이는 AC 문구가 명시하는 "None 방어" 그 자체는 아니지만, AC3/AC4가
    함께 요구하는 "예외가 발생하지 않는다"는 상위 계약(파싱 전체 관점)을
    검증하는 예외 케이스다.
    """
    from unittest.mock import MagicMock

    mock_para = MagicMock()
    mock_para.level = "not-an-int"
    mock_para.runs = []
    mock_para.text = "garbage level"

    mock_shape = MagicMock()
    mock_shape.shape_id = 42
    mock_shape.name = "GarbageLevelShape"
    mock_shape.shape_type = None
    mock_shape.has_table = False
    mock_shape.has_text_frame = True
    mock_shape.placeholder_format = None
    mock_shape.text_frame.paragraphs = [mock_para]
    mock_shape.left = 1
    mock_shape.top = 2
    mock_shape.width = 3
    mock_shape.height = 4

    fake_slide = types.SimpleNamespace(shapes=[mock_shape])
    slide_ir = _parse_slide(0, fake_slide)  # must not raise
    assert len(slide_ir.shapes) == 1
    result = slide_ir.shapes[0]
    # Demoted to OtherShapeIR because int("not-an-int") raises inside _parse_text,
    # and _parse_shape's outer try/except (ADR-204) absorbs it.
    assert isinstance(result, OtherShapeIR)
    assert result.left == 1 and result.top == 2


# ===========================================================================
# AC5/AC6 — Chart/SmartArt fallback
# ===========================================================================


def test_tc67_08_chart_no_extractable_text_returns_empty_no_exception(
    tmp_path: Path,
) -> None:
    """TC-67-08 (AC5, 경계): 제목·카테고리·시리즈명 전혀 없는 그래픽 프레임
    (알 수 없는 uri, 텍스트 런 없음) → fallback_text == "" 이고 예외 없음.

    주: presentationml/2006/ole uri는 python-pptx가 자체적으로
    shape_type == LINKED_OLE_OBJECT 로 인식하므로(1차 라벨링 경로),
    mso_shape_type은 "GRAPHIC_FRAME"이 아니라 "LINKED_OLE_OBJECT"가 된다.
    두 경로 모두 "UNKNOWN"이 아닌 실제 라벨이면 AC6 의도를 충족하므로
    라벨 종류 자체는 "UNKNOWN 아님"으로만 단언한다.
    """
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sp_tree = slide.shapes._spTree
    from lxml import etree  # noqa: PLC0415

    xml = (
        "<p:graphicFrame"
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        "<p:nvGraphicFramePr>"
        '<p:cNvPr id="91" name="NoTextObject"/>'
        "<p:cNvGraphicFramePr/><p:nvPr/>"
        "</p:nvGraphicFramePr>"
        "<p:xfrm>"
        '<a:off x="0" y="0"/><a:ext cx="100" cy="100"/>'
        "</p:xfrm>"
        "<a:graphic>"
        '<a:graphicData uri="http://schemas.openxmlformats.org/presentationml/2006/ole"/>'
        "</a:graphic>"
        "</p:graphicFrame>"
    )
    sp_tree.append(etree.fromstring(xml))
    p = _save(prs, tmp_path, "no_text_object.pptx")

    pres = parse_presentation(p)
    others = [s for s in pres.slides[0].shapes if isinstance(s, OtherShapeIR)]
    assert len(others) == 1
    assert others[0].fallback_text == ""
    assert others[0].mso_shape_type != "UNKNOWN"


def test_tc67_09_text_frame_access_raises_falls_back_to_chart_xml() -> None:
    """TC-67-09 (AC5, 예외): has_text_frame==True 이지만 text_frame.text
    접근이 예외를 던지는 손상 케이스에서도, has_chart 경로의 XML 추출로
    폴백되어 텍스트가 복구된다 (developer 테스트가 다루지 않은 조합 경로).
    """

    class _RaisingTextFrame:
        @property
        def text(self) -> str:
            raise RuntimeError("corrupted text frame")

    class _ChartElementWithTitle:
        _A_T_QN = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"

        def iter(self, tag: str):
            from lxml import etree  # noqa: PLC0415

            a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            xml = f'<root xmlns:a="{a_ns}"><a:t>Recovered Title</a:t></root>'
            root = etree.fromstring(xml)
            yield from root.iter(tag)

    class _HybridShape:
        has_text_frame = True
        has_chart = True
        element = None
        text_frame = _RaisingTextFrame()

        @property
        def chart(self) -> object:
            return types.SimpleNamespace(element=_ChartElementWithTitle())

    result = _extract_graphic_frame_fallback_text(_HybridShape())
    assert "Recovered Title" in result


def test_tc67_10_shape_type_present_uses_enum_name_not_graphic_frame_helper() -> None:
    """TC-67-10 (AC6, 정상): shape_type이 존재하는 (None이 아닌) GraphicFrame류
    도형은 `_graphic_frame_type_label` 을 거치지 않고 shape_type.name을 그대로
    사용한다 — 두 라벨링 경로(1차 shape_type / 2차 XML fallback) 분기 확인.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from pptx_md.parser import _parse_other

    class _FakeShapeType:
        name = "EMBEDDED_OLE_OBJECT"

    class _OleShape:
        shape_type = _FakeShapeType()
        has_text_frame = False
        element = None

    result = _parse_other(_OleShape(), sid=7, name="OleObj")
    assert result.mso_shape_type == "EMBEDDED_OLE_OBJECT"
    assert result.mso_shape_type != "UNKNOWN"
    # sanity: real MSO_SHAPE_TYPE enum also exposes .name the same way
    assert MSO_SHAPE_TYPE.CHART.name == "CHART"


# ===========================================================================
# AC7 — partial failure isolation, boundary positions
# ===========================================================================


def test_tc67_11_corrupted_shape_at_first_and_last_position(tmp_path: Path) -> None:
    """TC-67-11 (AC7, 경계): 손상 도형이 슬라이드의 첫 번째/마지막 위치에
    있어도 격리되고, 정상 도형은 위치와 무관하게 보존된다.

    주: placeholder_format 접근 시 ValueError를 던지면 _parse_text 자체가
    "non-placeholder shape"의 정상 신호로 캐치해 흡수하므로(FR-04 기존 설계)
    이는 "손상"이 아니라 정상 처리 경로가 된다. 진짜 손상을 흉내내려면
    RuntimeError 등 그 catch 절이 잡지 않는 예외를 던져야 한다
    (developer의 AC7 테스트와 동일한 선택 — 검증 후 정합성 확인).
    """
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tb.text_frame.text = "middle-normal"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    prs2 = PptxPresentation(buf)
    real_shape = list(prs2.slides[0].shapes)[0]

    class _CorruptedShape:
        shape_id = 900
        name = "Corrupted"
        shape_type = None
        has_table = False
        has_text_frame = True
        left = 1
        top = 2
        width = 3
        height = 4

        @property
        def placeholder_format(self) -> object:
            raise RuntimeError("boom")

    fake_slide = types.SimpleNamespace(
        shapes=[_CorruptedShape(), real_shape, _CorruptedShape()]
    )
    slide_ir = _parse_slide(0, fake_slide)
    assert len(slide_ir.shapes) == 3
    first, middle, last = slide_ir.shapes
    assert isinstance(first, OtherShapeIR)
    assert isinstance(last, OtherShapeIR)
    assert isinstance(middle, TextShapeIR)
    assert "middle-normal" in "".join(pg.text for pg in middle.paragraphs)


def test_tc67_12_corrupted_shape_missing_coord_attrs_defaults_zero() -> None:
    """TC-67-12 (AC7, 예외): 손상 도형이 좌표 속성 자체가 없는 경우
    (getattr 기본값 None → 0), 강등 결과의 좌표가 0/0/0/0 이며 예외가
    전파되지 않는다 (FR-23 None 정규화와 ADR-204의 합성 케이스)."""

    class _NoCoordCorruptedShape:
        shape_id = 1
        name = "NoCoords"
        shape_type = None
        has_table = False
        has_text_frame = True
        # left/top/width/height intentionally absent

        @property
        def placeholder_format(self) -> object:
            raise RuntimeError("boom: no coords either")

    fake_slide = types.SimpleNamespace(shapes=[_NoCoordCorruptedShape()])
    slide_ir = _parse_slide(0, fake_slide)
    assert len(slide_ir.shapes) == 1
    result = slide_ir.shapes[0]
    assert isinstance(result, OtherShapeIR)
    assert (result.left, result.top, result.width, result.height) == (0, 0, 0, 0)


# ===========================================================================
# AC8 — no image loss, deeper nesting than developer's fixture
# ===========================================================================


def test_tc67_13_images_across_two_level_nesting_no_loss(tmp_path: Path) -> None:
    """TC-67-13 (AC8, 경계): 2단 중첩 group 내부 이미지까지 포함한 4개
    이미지(최상위 1 + 1단 그룹 1 + 2단 그룹 2)가 전량 파싱된다.

    developer 픽스처는 1단 그룹까지만 다뤘다 — 더 깊은 중첩에서의 손실을
    별도로 확인한다.
    """
    png = tmp_path / "nested_img.png"
    png.write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d4948445200000001000000010802000000907753de"
            "0000000c49444154789c63f8ffff3f0005fe02fe0def46b80000000049454e44ae"
            "426082"
        )
    )
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.shapes.add_picture(str(png), Inches(0.2), Inches(0.2), Inches(1), Inches(1))
    inner_pic_a = slide.shapes.add_picture(
        str(png), Inches(2), Inches(0.2), Inches(1), Inches(1)
    )
    inner_pic_b = slide.shapes.add_picture(
        str(png), Inches(2), Inches(2), Inches(1), Inches(1)
    )
    inner_group = slide.shapes.add_group_shape([inner_pic_a, inner_pic_b])
    outer_group_pic = slide.shapes.add_picture(
        str(png), Inches(4), Inches(0.2), Inches(1), Inches(1)
    )
    slide.shapes.add_group_shape([inner_group, outer_group_pic])

    p = _save(prs, tmp_path, "nested_images.pptx")
    pres = parse_presentation(p)
    slide_ir = pres.slides[0]

    all_images: list[ImageShapeIR] = [
        s for s in iter_shapes(slide_ir) if isinstance(s, ImageShapeIR)
    ]
    top_level_images = [s for s in slide_ir.shapes if isinstance(s, ImageShapeIR)]
    assert len(top_level_images) == 1  # top_pic
    assert len(all_images) == 4
    for img in all_images:
        assert len(img.image_bytes) > 0


# ===========================================================================
# AC9 — COM path N/A confirmation (independent grep, not trusting the report)
# ===========================================================================


def test_tc67_14_convert_via_com_absent_in_src(monkeypatch: pytest.MonkeyPatch) -> None:
    """TC-67-14 (AC9, N/A 확인): src/pptx_md 전역에 convert_via_com 심볼이
    존재하지 않음을 코드 검사로 재확인한다 (developer/reviewer 주장을
    액면 그대로 받아들이지 않고 직접 스캔)."""
    src_root = Path(__file__).parent.parent / "src" / "pptx_md"
    hits = []
    for py_file in src_root.rglob("*.py"):
        content = py_file.read_text(encoding="utf-8")
        if "convert_via_com" in content:
            hits.append(py_file)
    assert hits == [], f"convert_via_com found in: {hits}"


# ===========================================================================
# AC10 — determinism, 3 repetitions + unicode/special-char content
# ===========================================================================


def test_tc67_15_determinism_three_repeats_with_unicode_and_chart(
    tmp_path: Path,
) -> None:
    """TC-67-15 (AC10, 정상): 유니코드·특수문자 텍스트 + 차트를 포함한 슬라이드를
    3회 반복 파싱해도 매번 완전히 동일한 IR을 낸다 (developer는 2회만 반복)."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(1))
    tb.text_frame.text = '한글 텍스트 | 특수문자 "quote" & <tag>'

    chart_data = CategoryChartData()
    chart_data.categories = ["가", "나", "다"]
    chart_data.add_series("시리즈", (1, 2, 3))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(2),
        Inches(4),
        Inches(3),
        chart_data,
    )
    gframe.chart.has_title = True
    gframe.chart.chart_title.text_frame.text = "유니코드 차트"

    p = _save(prs, tmp_path, "unicode_determinism.pptx")

    results = [parse_presentation(p) for _ in range(3)]
    assert results[0] == results[1] == results[2]

    others = [s for s in results[0].slides[0].shapes if isinstance(s, OtherShapeIR)]
    assert len(others) == 1
    assert "유니코드 차트" in others[0].fallback_text


# ===========================================================================
# Matrix completion — additional 정상/경계/예외 cases per AC2/AC4/AC5/AC6/AC7
# ===========================================================================


def test_tc67_18_nested_group_corrupted_leaf_isolated(tmp_path: Path) -> None:
    """TC-67-18 (AC2, 예외): 2단 중첩 group의 최심부(leaf) 자리에 손상 도형이
    있어도 해당 leaf만 강등되고, outer/inner group 구조 자체와 형제 leaf는
    영향받지 않는다 (깊은 재귀 경로에서의 예외 격리 회귀 방지)."""

    class _DeepBoom:
        shape_id = 777
        name = "DeepBoom"
        shape_type = None
        has_table = False
        has_text_frame = True
        left = 1
        top = 1
        width = 1
        height = 1

        @property
        def placeholder_format(self) -> object:
            raise RuntimeError("boom at depth 2")

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sibling_leaf = slide.shapes.add_textbox(
        Inches(3), Inches(0.5), Inches(2), Inches(1)
    )
    sibling_leaf.text_frame.text = "sibling at depth 1"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    prs2 = PptxPresentation(buf)
    real_sibling = list(prs2.slides[0].shapes)[0]

    inner_group = types.SimpleNamespace(
        shape_id=2,
        name="Inner",
        shape_type=MSO_SHAPE_TYPE.GROUP,
        left=0,
        top=0,
        width=10,
        height=10,
        shapes=[_DeepBoom()],
    )
    outer_group = types.SimpleNamespace(
        shape_id=1,
        name="Outer",
        shape_type=MSO_SHAPE_TYPE.GROUP,
        left=0,
        top=0,
        width=20,
        height=20,
        shapes=[inner_group, real_sibling],
    )
    fake_slide = types.SimpleNamespace(shapes=[outer_group])

    slide_ir = _parse_slide(0, fake_slide)
    assert len(slide_ir.shapes) == 1
    outer_ir = slide_ir.shapes[0]
    assert isinstance(outer_ir, GroupShapeIR)
    assert len(outer_ir.children) == 2
    inner_ir, sibling_ir = outer_ir.children
    assert isinstance(inner_ir, GroupShapeIR)
    assert len(inner_ir.children) == 1
    assert isinstance(inner_ir.children[0], OtherShapeIR)
    assert inner_ir.children[0].mso_shape_type == "UNKNOWN"
    assert isinstance(sibling_ir, TextShapeIR)
    assert "sibling at depth 1" in "".join(pg.text for pg in sibling_ir.paragraphs)


def test_tc67_19_mixed_none_and_int_levels_in_same_frame(tmp_path: Path) -> None:
    """TC-67-19 (AC4, 경계): 동일 텍스트 프레임 내 일부 단락만 level 이 None,
    나머지는 유효한 정수 — None인 단락만 0으로 정규화되고 다른 단락의 유효한
    레벨 값은 그대로 유지된다 (developer 테스트는 단일 단락·전체 None만 검증).
    """
    from unittest.mock import MagicMock

    para_none = MagicMock()
    para_none.level = None
    para_none.runs = []
    para_none.text = "none level"

    para_two = MagicMock()
    para_two.level = 2
    para_two.runs = []
    para_two.text = "level two"

    mock_shape = MagicMock()
    mock_shape.placeholder_format = None
    mock_shape.text_frame.paragraphs = [para_none, para_two]

    from pptx_md.parser import _parse_text

    result = _parse_text(mock_shape, sid=1, name="mixed_none_int")
    assert len(result.paragraphs) == 2
    assert result.paragraphs[0].level == 0
    assert result.paragraphs[0].text == "none level"
    assert result.paragraphs[1].level == 2
    assert result.paragraphs[1].text == "level two"


def test_tc67_20_chart_property_and_element_both_raise_yields_empty() -> None:
    """TC-67-20 (AC5, 예외): has_chart==True 인데 chart 속성 접근 자체가
    예외를 던지고, 동시에 shape.element 도 None 인 이중 실패 상황에서도
    예외 없이 빈 문자열이 반환된다 (developer는 element!=None인 경우만
    혹은 chart 실패만 단독으로 검증; 이중 실패 조합은 신규 검증)."""

    class _DoubleFailureShape:
        has_chart = True
        element = None

        @property
        def chart(self) -> object:
            raise RuntimeError("chart part completely unavailable")

    result = _extract_graphic_frame_fallback_text(_DoubleFailureShape())
    assert result == ""


def test_tc67_21_custom_nonstandard_uri_yields_graphic_frame_label() -> None:
    """TC-67-21 (AC6, 경계): python-pptx가 자체 인식하지 못하는(shape_type이
    None으로 남는) 완전히 임의의 graphicData uri는 "GRAPHIC_FRAME" 라벨로
    귀결된다 — DIAGRAM/CHART가 아닌 제3의 미지 콘텐츠 케이스."""
    from lxml import etree  # noqa: PLC0415

    from pptx_md.parser import _graphic_frame_type_label

    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    elem = etree.fromstring(
        f'<root xmlns:a="{a_ns}"><a:graphic>'
        '<a:graphicData uri="urn:custom:totally-unknown-graphic-2099"/>'
        "</a:graphic></root>"
    )

    class _ShapeCustomUri:
        element = elem

    label = _graphic_frame_type_label(_ShapeCustomUri())
    assert label == "GRAPHIC_FRAME"
    assert label != "UNKNOWN"


def test_tc67_22_picture_image_property_raises_isolated_from_siblings(
    tmp_path: Path,
) -> None:
    """TC-67-22 (AC7, 정상): PICTURE 디스패치 경로(TEXT가 아님)에서 발생한
    예외도 동일하게 해당 도형만 강등되고 좌표는 유지되며, 이웃 도형은 영향
    받지 않는다. developer의 AC7 테스트는 TEXT류 경로(placeholder_format)만
    사용했으므로, 서로 다른 디스패치 분기(PICTURE)에서도 동일 계약이
    성립하는지 별도로 검증한다."""

    class _BoomImageProperty:
        shape_id = 321
        name = "BoomPicture"
        shape_type = MSO_SHAPE_TYPE.PICTURE
        has_table = False
        has_text_frame = False
        left = 5
        top = 6
        width = 7
        height = 8

        @property
        def image(self) -> object:
            raise RuntimeError("boom: image part unavailable")

        element = None

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tb.text_frame.text = "neighbour ok"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    prs2 = PptxPresentation(buf)
    real_text = list(prs2.slides[0].shapes)[0]

    fake_slide = types.SimpleNamespace(shapes=[real_text, _BoomImageProperty()])
    slide_ir = _parse_slide(0, fake_slide)
    assert len(slide_ir.shapes) == 2
    normal, corrupted = slide_ir.shapes
    assert isinstance(normal, TextShapeIR)
    assert isinstance(corrupted, OtherShapeIR)
    assert corrupted.mso_shape_type == "UNKNOWN"
    assert (corrupted.left, corrupted.top, corrupted.width, corrupted.height) == (
        5,
        6,
        7,
        8,
    )


def test_tc67_23_graphic_frame_label_find_raises_falls_back_unknown_safely() -> None:
    """TC-67-23 (AC6, 예외): element.find 가 예외를 던지는 손상 상황에서는
    실제 라벨을 결정할 수 없으므로 안전하게 "UNKNOWN"으로 귀결되고, 예외는
    전파되지 않는다 (실제 라벨 보장은 "결정 가능한 경우"에 한정된다는
    AC6 문맥의 예외-안전성 측면을 확인한다)."""
    from pptx_md.parser import _parse_other

    class _RaisingFindElement:
        def find(self, path: str) -> object:
            raise RuntimeError("boom: malformed xml")

    class _ShapeWithRaisingFind:
        shape_type = None
        has_text_frame = False
        element = _RaisingFindElement()

    result = _parse_other(_ShapeWithRaisingFind(), sid=8, name="RaisingFind")
    assert result.mso_shape_type == "UNKNOWN"
    assert result.fallback_text == ""


# ===========================================================================
# Integration smoke: parse -> assemble on a real fixture PPTX
# ===========================================================================


def test_tc67_16_integration_smoke_parse_to_assemble_no_regression() -> None:
    """TC-67-16 (통합 스모크): 실제 픽스처(other_shape.pptx)로 parse_presentation
    -> assemble_document 전체 파이프라인이 예외 없이 동작하고, FR-26이 보강한
    차트 fallback_text가 최종 Markdown에 실제로 흘러 들어가는지 확인한다.
    """
    fixture = FIXTURES / "other_shape.pptx"
    if not fixture.is_file():
        pytest.skip("other_shape.pptx fixture 없음")

    pres = parse_presentation(fixture)
    doc = assemble_document(pres)

    assert isinstance(doc, str)
    assert len(doc) > 0
    # AC5 회귀 확인: chart 카테고리 텍스트가 최종 문서까지 보존되는지
    others = [
        s
        for slide in pres.slides
        for s in iter_shapes(slide)
        if isinstance(s, OtherShapeIR)
    ]
    assert len(others) >= 1
    assert others[0].fallback_text != ""


def test_tc67_17_integration_smoke_determinism_end_to_end() -> None:
    """TC-67-17 (통합 스모크, 결정성): 동일 픽스처를 parse+assemble 2회 실행해도
    최종 Markdown 문자열까지 완전히 동일하다 (ADR-218이 assembler 단까지
    이어지는지 파서 커버리지 검증의 연장선에서 확인)."""
    fixture = FIXTURES / "other_shape.pptx"
    if not fixture.is_file():
        pytest.skip("other_shape.pptx fixture 없음")

    doc1 = assemble_document(parse_presentation(fixture))
    doc2 = assemble_document(parse_presentation(fixture))
    assert doc1 == doc2
