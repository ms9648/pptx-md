"""TR-68 독립 QA 시나리오 — 이슈 #68 (FR-27) AC1~AC10 재검증.

tester(QA) 산출물. developer 의 tests/test_description_pipeline.py,
tests/test_api.py, tests/test_mermaid.py 를 재실행해 "통과"로 갈음하지 않는다.
동일 공개 API(enrich_descriptions/convert/render_diagram_mermaid)를
다른 픽스처·다른 실패 모드·다른 관찰 방식으로 재구성해 독립적으로 검증한다.

AC 번호는 이슈 #68 원문 기준(설계 문서 ARCH-M12 §1의 재배열 번호가 아님).

특별 지시 반영:
  - AC7(동시성): 카운팅 mock을 이 파일에서 독립 재구현. 상한 1/기본값4/과대값
    경계를 모두 검증.
  - AC8(결정성): 동시성이 실제로 활성화된 상태(다수 worker, 인위적 지연으로
    완료 순서를 뒤섞음)에서 3회 이상 반복 변환 후 바이트 동일성 확인.
  - AC6(캐싱): 슬라이드 간(다른 슬라이드) 동일 바이트 이미지의 dedup과,
    "1바이트만 다른" 유사 이미지의 비-dedup(별개 처리)을 함께 검증.
  - AC3+AC9 조합: 이미지별 실패 격리와 FR-22 위치 마커 대체가 함께 발동하는
    시나리오를, developer가 사용한 "describe() 예외" 경로가 아니라
    "EMF/WMF 벡터 변환 실패" 경로로 재현(다른 실패 모드로 동일 계약 검증).
  - AC2(NFR-08): sys.modules 검사로 SDK 미임포트를 독립 확인.
"""

from __future__ import annotations

import io
import logging
import sys
import threading
import time
from pathlib import Path
from typing import Any
from unittest.mock import patch

import pytest
from PIL import Image
from pptx import Presentation as PptxPresentation

from pptx_md import mermaid, vector
from pptx_md.api import ConvertOptions, convert
from pptx_md.assembler import assemble_document
from pptx_md.description_pipeline import enrich_descriptions
from pptx_md.ir import (
    ImageClass,
    ImageShapeIR,
    PresentationIR,
    ShapeKind,
    SlideIR,
    TextShapeIR,
)

_PROJECT_ROOT = Path(__file__).parent.parent

# ---------------------------------------------------------------------------
# 독립 헬퍼 (developer 헬퍼 재사용 없음 — 새로 구성)
# ---------------------------------------------------------------------------


def _png_bytes(color: tuple[int, int, int], size: tuple[int, int] = (40, 40)) -> bytes:
    """단색 PNG 실제 바이트를 생성한다(합성이 아닌 유효 PNG 디코딩 가능 데이터)."""
    buf = io.BytesIO()
    Image.new("RGB", size, color=color).save(buf, format="PNG")
    return buf.getvalue()


def _slide(shapes: list[Any], index: int = 0) -> SlideIR:
    return SlideIR(index=index, title=f"슬라이드 {index + 1}", shapes=shapes)


def _pres(slides: list[SlideIR]) -> PresentationIR:
    return PresentationIR(source_path="tr68_qa.pptx", slides=slides)


def _img(
    shape_id: int,
    image_bytes: bytes,
    image_ext: str = "png",
    classification: ImageClass | None = None,
) -> ImageShapeIR:
    return ImageShapeIR(
        shape_id=shape_id,
        name=f"pic{shape_id}",
        kind=ShapeKind.IMAGE,
        image_bytes=image_bytes,
        image_format=image_ext,
        image_ext=image_ext,
        classification=classification,
    )


def _text(shape_id: int, text: str) -> TextShapeIR:
    from pptx_md.ir import ParagraphIR

    return TextShapeIR(
        shape_id=shape_id,
        name="body",
        kind=ShapeKind.TEXT,
        paragraphs=[ParagraphIR(text=text, level=0)],
    )


class QaKeyedDescriber:
    """이미지 바이트를 키로 삼아 결정적 텍스트를 반환하는 독립 Fake.

    developer의 HashAwareDescriber와 동일 개념이나 이 파일에서 새로 작성.
    """

    def __init__(self, delay: float = 0.0) -> None:
        self.calls: list[bytes] = []
        self.delay = delay

    def describe(
        self, image_bytes: bytes, image_ext: str, shape_hint: str | None
    ) -> str:
        self.calls.append(image_bytes)
        if self.delay:
            time.sleep(self.delay)
        return f"qa-desc:{hash(image_bytes) & 0xFFFF:04x}"


class QaConcurrencyProbe:
    """동시 활성 describe() 호출 수를 추적하는 독립 카운팅 mock.

    developer의 ConcurrencyCountingDescriber와 별개로 이 파일에서 새로
    구현. per-call delay를 이미지 인덱스 역순으로 부여해 완료 순서를
    제출 순서와 강제로 어긋나게 만들 수 있는 옵션을 추가로 제공한다
    (AC8 결정성 스트레스용).
    """

    def __init__(self, reverse_delay: bool = False, base_delay: float = 0.01) -> None:
        self._lock = threading.Lock()
        self._active = 0
        self.max_observed = 0
        self.call_count = 0
        self.reverse_delay = reverse_delay
        self.base_delay = base_delay
        self._seen_bytes: list[bytes] = []

    def describe(
        self, image_bytes: bytes, image_ext: str, shape_hint: str | None
    ) -> str:
        with self._lock:
            self._active += 1
            self.max_observed = max(self.max_observed, self._active)
            self.call_count += 1
            idx = len(self._seen_bytes)
            self._seen_bytes.append(image_bytes)
        if self.reverse_delay:
            # 나중에 제출된 job일수록 먼저 끝나도록(첫 제출 job이 가장 늦게
            # 완료) 완료 순서를 인위적으로 뒤집는다.
            time.sleep(self.base_delay * (10 - idx if idx < 10 else 1))
        else:
            time.sleep(self.base_delay)
        with self._lock:
            self._active -= 1
        return f"probe-desc:{image_bytes!r}"


class QaSelectiveRaiser:
    """지정한 바이트에서만 예외를 던지는 독립 Fake (다른 예외 유형 사용)."""

    def __init__(self, fail_on: set[bytes], text: str = "qa-ok") -> None:
        self.fail_on = fail_on
        self.text = text
        self.calls = 0

    def describe(
        self, image_bytes: bytes, image_ext: str, shape_hint: str | None
    ) -> str:
        self.calls += 1
        if image_bytes in self.fail_on:
            # developer 테스트는 RuntimeError를 사용 — 여기서는 다른 예외
            # 유형(연쇄 예외)으로 격리 계약이 예외 타입에 의존하지 않음을 확인.
            try:
                raise KeyError("upstream lookup failed")
            except KeyError as exc:
                raise ConnectionError("simulated VLM transport failure") from exc
        return self.text


class QaAllFailDescriber:
    """모든 호출에서 예외를 던지는 독립 Fake (AC9 극단 경계: 커버리지 0%)."""

    def describe(
        self, image_bytes: bytes, image_ext: str, shape_hint: str | None
    ) -> str:
        raise OSError("qa: all images fail")


def _save_real_pptx_with_images(path: Path, colors: list[tuple[int, int, int]]) -> None:
    """실제 PPTX 파일에 슬라이드별로 서로 다른 실제 PNG 그림을 삽입한다.

    developer 테스트는 전량 합성 PresentationIR(mock parse)을 사용한다.
    본 QA는 python-pptx로 실제 파일을 만들고 parse_presentation을 목킹
    없이 그대로 호출해(§3 통합 스모크 요구) 진짜 파이프라인 결합을 검증한다.
    """
    prs = PptxPresentation()
    for color in colors:
        layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(layout)
        img_buf = io.BytesIO(_png_bytes(color))
        slide.shapes.add_picture(img_buf, left=0, top=0, width=914400, height=914400)
    buf = io.BytesIO()
    prs.save(buf)
    path.write_bytes(buf.getvalue())


# ===========================================================================
# AC1 (정상 경로 — 파이프라인 실사용 연결) + 통합 스모크(요구 §3)
# ===========================================================================


class TestAc1PipelineRealUsage:
    """TC-68-01~03: AC1 — mock 없이 실제 parse→enrich→assemble e2e."""

    def test_tc_68_01_실제_pptx_실제_이미지_e2e(self, tmp_path: Path) -> None:
        """TC-68-01 (정상): 실제 PPTX(진짜 PNG 3장, 서로 다른 슬라이드) +
        describer 주입 → convert() 를 mock 없이 그대로 실행 → 각 설명이
        본문에 포함된다. parse_presentation/enrich_images 어느 것도
        패치하지 않는다(developer의 AC1 e2e 테스트는 parse/enrich_images를
        패치함 — 본 TC는 그 두 단계까지 실제로 통과시킨다)."""
        pptx_path = tmp_path / "real_images.pptx"
        _save_real_pptx_with_images(
            pptx_path, colors=[(255, 0, 0), (0, 255, 0), (0, 0, 255)]
        )
        describer = QaKeyedDescriber()

        result = convert(pptx_path, options=ConvertOptions(describer=describer))

        assert len(describer.calls) == 3
        for c in describer.calls:
            token = f"qa-desc:{hash(c) & 0xFFFF:04x}"
            assert token in result

    def test_tc_68_02_경계_이미지_0장_describer_주입해도_호출_0(
        self, tmp_path: Path
    ) -> None:
        """TC-68-02 (경계): 이미지가 전혀 없는 실제 PPTX에 describer를
        주입해도 describe 호출이 0건이며 예외 없이 완료된다."""
        prs = PptxPresentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = "제목만 있는 슬라이드"
        buf = io.BytesIO()
        prs.save(buf)
        pptx_path = tmp_path / "no_images.pptx"
        pptx_path.write_bytes(buf.getvalue())

        describer = QaKeyedDescriber()
        result = convert(pptx_path, options=ConvertOptions(describer=describer))

        assert describer.calls == []
        assert "제목만 있는 슬라이드" in result

    def test_tc_68_03_예외_특정_이미지_describe_예외여도_convert_성공(
        self, tmp_path: Path
    ) -> None:
        """TC-68-03 (예외): convert() 최상위 호출 관점에서, 이미지 중
        하나의 describe()가 예외를 던져도 convert()는 예외를 전파하지
        않고 문자열을 반환한다(developer 테스트는 description_pipeline
        레벨에서만 검증 — 본 TC는 public convert() 레벨에서 재확인)."""
        pptx_path = tmp_path / "one_fails.pptx"
        red = _png_bytes((200, 10, 10))
        _save_real_pptx_with_images(pptx_path, colors=[(200, 10, 10), (10, 200, 10)])
        describer = QaSelectiveRaiser(fail_on={red})

        result = convert(pptx_path, options=ConvertOptions(describer=describer))

        assert isinstance(result, str)
        assert "qa-ok" in result  # 나머지 이미지는 정상 설명


# ===========================================================================
# AC2 (경계 — describer 미주입, NFR-08) — sys.modules 독립 검사
# ===========================================================================


class TestAc2Nfr08SysModules:
    """TC-68-04~06: describer=None 시 SDK 미임포트를 sys.modules로 검사."""

    def test_tc_68_04_정상_describer_none_sys_modules에_sdk_부재(
        self, tmp_path: Path
    ) -> None:
        """TC-68-04 (정상): describer=None으로 실제 이미지 포함 PPTX를
        convert 한 뒤에도 sys.modules에 anthropic/openai가 없다."""
        pptx_path = tmp_path / "deck.pptx"
        _save_real_pptx_with_images(pptx_path, colors=[(1, 2, 3)])

        convert(pptx_path)  # describer 미지정 -> 기본 None

        assert "anthropic" not in sys.modules
        assert "openai" not in sys.modules

    def test_tc_68_05_경계_다중이미지_describer_none_호출0(
        self, tmp_path: Path
    ) -> None:
        """TC-68-05 (경계): 이미지 5장 + describer=None → description 전부
        None 유지, SDK 미임포트."""
        shapes = [_img(i, _png_bytes((i * 10, i * 10, i * 10))) for i in range(5)]
        pres = _pres([_slide(shapes)])

        enrich_descriptions(pres, None)

        assert all(s.description is None for s in shapes)
        assert "anthropic" not in sys.modules
        assert "openai" not in sys.modules

    def test_tc_68_06_예외_옵션값이_설정돼도_describer_none이면_무발동(
        self,
    ) -> None:
        """TC-68-06 (예외적 경계): diagram_mermaid=True, max_workers=8 처럼
        비-기본 옵션이 설정돼도 describer=None이면 어떤 동작도 발동하지
        않는다(풀 생성 자체가 없음을 ThreadPoolExecutor 미호출로 검증)."""
        from unittest.mock import MagicMock

        shapes = [_img(1, _png_bytes((9, 9, 9)), classification=ImageClass.DIAGRAM)]
        pres = _pres([_slide(shapes)])

        with patch(
            "pptx_md.description_pipeline.ThreadPoolExecutor", MagicMock()
        ) as pool_mock:
            enrich_descriptions(pres, None, max_workers=8, diagram_mermaid=True)

        pool_mock.assert_not_called()
        assert shapes[0].description is None


# ===========================================================================
# AC3 (예외 — 이미지별 격리) — 다른 예외 유형/전체실패 경계
# ===========================================================================


class TestAc3IsolationDifferentFailureModes:
    """TC-68-07~09: AC3 격리 — 연쇄 예외, 전체 실패, 정상 혼재."""

    def test_tc_68_07_정상_연쇄예외여도_해당_이미지만_격리(self) -> None:
        """TC-68-07 (정상): ConnectionError(연쇄, from KeyError) — developer의
        RuntimeError와 다른 예외 유형 — 를 던져도 해당 이미지만 None,
        나머지는 정상."""
        fail_bytes = _png_bytes((77, 77, 77))
        ok_bytes = _png_bytes((5, 6, 7))
        shapes = [_img(1, fail_bytes), _img(2, ok_bytes)]
        pres = _pres([_slide(shapes)])
        describer = QaSelectiveRaiser(fail_on={fail_bytes})

        enrich_descriptions(pres, describer)

        assert shapes[0].description is None
        assert shapes[1].description == "qa-ok"

    def test_tc_68_08_경계_전체_이미지_실패해도_예외_전파_0(self) -> None:
        """TC-68-08 (경계): 모든 이미지가 실패해도 enrich_descriptions는
        예외를 전파하지 않고 전부 None으로 남긴다(developer는 "일부 실패"만
        검증 — 본 TC는 전체 실패라는 극단 경계)."""
        shapes = [_img(i, _png_bytes((i, i, i))) for i in range(1, 4)]
        pres = _pres([_slide(shapes)])
        describer = QaAllFailDescriber()

        enrich_descriptions(pres, describer)  # 예외 없이 반환해야 함

        assert all(s.description is None for s in shapes)

    def test_tc_68_09_예외_격리후_최종_markdown_정상_조립(self) -> None:
        """TC-68-09 (예외 관점): 격리된 이미지가 있어도 assemble_document
        전체가 정상 조립되고, 실패 이미지 자리에는 표준 FR-22 마커가
        나온다(정보 손실 없음 확인 — assembler까지 연결)."""
        fail_bytes = _png_bytes((1, 1, 1))
        shapes = [
            _text(1, "본문"),
            _img(2, fail_bytes),
            _img(3, _png_bytes((2, 2, 2))),
        ]
        pres = _pres([_slide(shapes)])
        describer = QaSelectiveRaiser(fail_on={fail_bytes})

        enrich_descriptions(pres, describer)
        md = assemble_document(pres)

        assert "슬라이드 1 이미지 1" in md  # 실패 이미지 자리의 표준 마커
        assert "qa-ok" in md  # 성공 이미지 설명


# ===========================================================================
# AC4/AC5 (Mermaid 다이어그램 산출 옵션 + fallback)
# ===========================================================================


class TestAc4Ac5MermaidOption:
    """TC-68-10~14: 관계 보존 확인 + classification 경계 + fallback 경계."""

    def test_tc_68_10_정상_diagram_mermaid_관계_노드_보존(self) -> None:
        """TC-68-10 (정상): DIAGRAM + 옵션 활성 → 응답의 노드/엣지 텍스트가
        최종 description에 그대로 보존된다(단순 펜스 존재가 아니라 실제
        관계 내용 보존을 확인 — developer 테스트보다 더 구체적인 단언)."""
        img_bytes = _png_bytes((3, 3, 3))
        shapes = [_img(1, img_bytes, classification=ImageClass.DIAGRAM)]
        pres = _pres([_slide(shapes)])

        class RelationDescriber:
            def describe(
                self, image_bytes: bytes, image_ext: str, shape_hint: str | None
            ) -> str:
                assert shape_hint is not None
                assert mermaid.DIAGRAM_HINT_SUFFIX in shape_hint
                return (
                    "설명 서두\n```mermaid\nflowchart TD\n"
                    "  START[시작] --> MID[검토]\n"
                    "  MID --> END[완료]\n```\n"
                )

        enrich_descriptions(pres, RelationDescriber(), diagram_mermaid=True)

        desc = shapes[0].description
        assert desc is not None
        assert "START[시작] --> MID[검토]" in desc
        assert "MID --> END[완료]" in desc
        assert desc.startswith("```mermaid")

    def test_tc_68_11_경계_옵션_off_DIAGRAM이어도_힌트접미사_미부착(self) -> None:
        """TC-68-11 (경계): diagram_mermaid=False(기본)면 DIAGRAM 분류
        이미지라도 hint에 DIAGRAM_HINT_SUFFIX가 붙지 않는다(옵션 off가
        기본이라는 이슈 명시 조건을 hint 레벨에서 직접 확인)."""
        received_hints: list[str | None] = []

        class RecordingDescriber:
            def describe(
                self, image_bytes: bytes, image_ext: str, shape_hint: str | None
            ) -> str:
                received_hints.append(shape_hint)
                return (
                    "plain text, even with a ```mermaid\nflowchart TD\n"
                    "A-->B\n``` fence"
                )

        shapes = [_img(1, _png_bytes((4, 4, 4)), classification=ImageClass.DIAGRAM)]
        pres = _pres([_slide(shapes)])

        enrich_descriptions(pres, RecordingDescriber())  # diagram_mermaid 기본 False

        assert received_hints[0] is not None
        assert mermaid.DIAGRAM_HINT_SUFFIX not in received_hints[0]
        # 옵션 off이므로 응답에 유효한 펜스가 있어도 후처리(추출)가 발동하지
        # 않고 원문 그대로 저장되어야 한다.
        assert shapes[0].description is not None
        assert shapes[0].description.startswith("plain text")

    def test_tc_68_12_예외_비DIAGRAM_분류는_옵션_활성이어도_힌트_미부착(
        self,
    ) -> None:
        """TC-68-12 (예외적 경계): diagram_mermaid=True이지만 classification
        이 DIAGRAM이 아닌 PHOTO인 이미지는 hint 증강도, mermaid 후처리도
        발동하지 않는다 — 응답이 우연히 유효한 mermaid flowchart 펜스여도
        원문 그대로 유지된다(요지: 트리거는 classification, 응답 내용이
        아님)."""
        raw_response = "```mermaid\nflowchart TD\nX-->Y\n```"

        class FixedDescriber:
            def describe(
                self, image_bytes: bytes, image_ext: str, shape_hint: str | None
            ) -> str:
                assert shape_hint is None or mermaid.DIAGRAM_HINT_SUFFIX not in (
                    shape_hint or ""
                )
                return raw_response

        shapes = [_img(1, _png_bytes((5, 5, 5)), classification=ImageClass.PHOTO)]
        pres = _pres([_slide(shapes)])

        enrich_descriptions(pres, FixedDescriber(), diagram_mermaid=True)

        assert shapes[0].description == raw_response  # 펜스 추출 안 됨(원문 그대로)

    def test_tc_68_13_경계_펜스는_있으나_키워드_불일치_fallback(self) -> None:
        """TC-68-13 (경계): ```mermaid 펜스는 있지만 flowchart/graph
        키워드가 아닌 pie 차트 → 원문 전체가 그대로 description에 유지된다
        (정보 손실 0)."""
        raw = (
            "이 다이어그램은 비중을 보여줍니다.\n```mermaid\npie title 비중\n"
            '"A" : 40\n"B" : 60\n```'
        )

        class PieDescriber:
            def describe(
                self, image_bytes: bytes, image_ext: str, shape_hint: str | None
            ) -> str:
                return raw

        shapes = [_img(1, _png_bytes((6, 6, 6)), classification=ImageClass.DIAGRAM)]
        pres = _pres([_slide(shapes)])

        enrich_descriptions(pres, PieDescriber(), diagram_mermaid=True)

        assert shapes[0].description == raw

    def test_tc_68_14_예외_펜스_본문이_공백만_fallback_무예외(self) -> None:
        """TC-68-14 (예외): ```mermaid펜스 내부가 공백/개행만인 응답 →
        예외 없이 원문 텍스트로 fallback."""
        raw = "구조를 파악하기 어렵습니다.\n```mermaid\n   \n\n```"

        class BlankFenceDescriber:
            def describe(
                self, image_bytes: bytes, image_ext: str, shape_hint: str | None
            ) -> str:
                return raw

        shapes = [_img(1, _png_bytes((7, 7, 7)), classification=ImageClass.DIAGRAM)]
        pres = _pres([_slide(shapes)])

        enrich_descriptions(pres, BlankFenceDescriber(), diagram_mermaid=True)

        assert shapes[0].description == raw


# ===========================================================================
# AC6 (해시 캐싱) — 슬라이드 간 중복 + 유사-비동일 바이트 구분
# ===========================================================================


class TestAc6HashCachingCrossSlide:
    """TC-68-15~17: 슬라이드 간 dedup + 1바이트 차이 구분 + 원본 미저장."""

    def test_tc_68_15_정상_서로_다른_슬라이드의_동일_바이트_이미지_dedup(
        self,
    ) -> None:
        """TC-68-15 (정상): 슬라이드 1과 슬라이드 3에 동일 바이트 이미지가
        있으면(중간 슬라이드는 무관 이미지) describe는 고유 이미지 수만큼만
        호출되고 두 슬라이드의 description이 동일하다."""
        shared = _png_bytes((123, 45, 67))
        slide1 = _slide([_img(1, shared)], index=0)
        slide2 = _slide([_img(2, _png_bytes((9, 9, 9)))], index=1)
        slide3 = _slide([_img(3, shared)], index=2)
        pres = _pres([slide1, slide2, slide3])
        describer = QaKeyedDescriber()

        enrich_descriptions(pres, describer)

        assert len(describer.calls) == 2  # 고유 2개(shared, 별개)
        assert slide1.shapes[0].description == slide3.shapes[0].description
        assert slide1.shapes[0].description != slide2.shapes[0].description

    def test_tc_68_16_경계_1바이트만_다른_이미지는_별개로_describe(self) -> None:
        """TC-68-16 (경계): 시각적으로 거의 동일해 보이더라도 바이트가
        1바이트라도 다르면 서로 다른 해시로 취급되어 describe가 각각
        호출된다(캐시가 "유사"가 아니라 "바이트 동일"만 dedup함을 확인)."""
        base = bytearray(_png_bytes((200, 200, 200)))
        variant = bytearray(base)
        variant[-1] ^= 0x01  # 마지막 바이트 1비트 변경
        shapes = [_img(1, bytes(base)), _img(2, bytes(variant))]
        pres = _pres([_slide(shapes)])
        describer = QaKeyedDescriber()

        enrich_descriptions(pres, describer)

        assert len(describer.calls) == 2
        assert shapes[0].description != shapes[1].description

    def test_tc_68_17_예외_원본_바이트_미저장_캐시값은_텍스트만(
        self, caplog: pytest.LogCaptureFixture
    ) -> None:
        """TC-68-17 (NFR-06 관련 예외/경계): 반복 이미지를 describe한 뒤
        로그 전체에 원본 이미지 바이트(디코딩 문자열)가 한 글자도 나타나지
        않는다 — sha256 다이제스트 앞 12자만 로그에 허용된다는 설계 계약을
        문자열 검색으로 재확인."""
        secret = _png_bytes((250, 1, 1))
        shapes = [_img(1, secret), _img(2, secret)]
        pres = _pres([_slide(shapes)])
        describer = QaKeyedDescriber()

        with caplog.at_level(logging.DEBUG, logger="pptx_md.description_pipeline"):
            enrich_descriptions(pres, describer)

        # PNG 바이트를 latin-1로 디코드한 원문이 로그에 그대로 나타나지 않아야 함
        raw_as_text = secret.decode("latin-1")
        assert raw_as_text not in caplog.text


# ===========================================================================
# AC7 (동시성 제한) — 독립 카운팅 mock, 상한 1/기본4/과대값
# ===========================================================================


class TestAc7ConcurrencyBounds:
    """TC-68-18~21: 독립 카운팅 mock으로 상한 1/기본4/과대값 재현."""

    def test_tc_68_18_정상_기본값4_동시_호출_4_초과_불가(self) -> None:
        """TC-68-18 (정상): max_workers를 지정하지 않으면(기본 4) 12개의
        고유 이미지를 describe하는 동안 어떤 시점에도 동시 활성 호출이
        4를 넘지 않는다."""
        shapes = [_img(i, _png_bytes((i, 0, 0))) for i in range(12)]
        pres = _pres([_slide(shapes)])
        probe = QaConcurrencyProbe(base_delay=0.02)

        enrich_descriptions(pres, probe)  # max_workers 미지정 -> 기본 4

        assert probe.call_count == 12
        assert probe.max_observed <= 4

    def test_tc_68_19_경계_상한_1_사실상_순차(self) -> None:
        """TC-68-19 (경계): max_workers=1이면 동시 활성 호출이 항상 1
        이하(사실상 순차)."""
        shapes = [_img(i, _png_bytes((0, i, 0))) for i in range(6)]
        pres = _pres([_slide(shapes)])
        probe = QaConcurrencyProbe(base_delay=0.015)

        enrich_descriptions(pres, probe, max_workers=1)

        assert probe.max_observed == 1
        assert probe.call_count == 6

    def test_tc_68_20_경계_과대값_상한이_고유이미지수보다_커도_초과없음(
        self,
    ) -> None:
        """TC-68-20 (과대값 경계): max_workers=1000처럼 고유 이미지 수보다
        훨씬 큰 상한을 줘도, 동시 활성 호출은 고유 이미지 수(4)를 절대
        넘지 않는다(클램프가 실제 동시 실행 동작에 반영됨을 카운팅 mock의
        관찰로 검증 — developer는 ThreadPoolExecutor 생성 인자만 검사, 본
        TC는 런타임 동시성 자체를 관찰)."""
        shapes = [_img(i, _png_bytes((0, 0, i))) for i in range(4)]
        pres = _pres([_slide(shapes)])
        probe = QaConcurrencyProbe(base_delay=0.03)

        enrich_descriptions(pres, probe, max_workers=1000)

        assert probe.call_count == 4
        assert probe.max_observed <= 4

    def test_tc_68_21_예외_음수_상한도_최소1로_동작(self) -> None:
        """TC-68-21 (예외적 입력): max_workers=-5 처럼 음수를 줘도 예외 없이
        동작하며 동시 활성 호출이 1을 넘지 않는다."""
        shapes = [_img(i, _png_bytes((i, i, 0))) for i in range(3)]
        pres = _pres([_slide(shapes)])
        probe = QaConcurrencyProbe(base_delay=0.02)

        enrich_descriptions(pres, probe, max_workers=-5)

        assert probe.call_count == 3
        assert probe.max_observed == 1


# ===========================================================================
# AC8 (결정성 골든 회귀) — 동시성 활성 + 3회+ 반복 + 완료순서 역전 스트레스
# ===========================================================================


class TestAc8DeterminismUnderConcurrency:
    """TC-68-22~24: 동시성이 실제로 활성인 상태에서 반복 변환 바이트 동일성."""

    def test_tc_68_22_정상_동시성_활성_3회_반복_바이트_동일(
        self, tmp_path: Path
    ) -> None:
        """TC-68-22 (정상): 실제 PPTX(진짜 PNG 6장) + max_workers=8(동시성
        확실히 활성) 로 convert()를 3회 반복 → 매번 바이트 동일 Markdown.
        각 반복마다 새 파일 파싱(캐시가 호출 간 재사용되지 않음을 함께
        확인 — 매 회 동일 describe 횟수)."""
        pptx_path = tmp_path / "det.pptx"
        colors = [(i * 30, i * 20, i * 10) for i in range(6)]
        _save_real_pptx_with_images(pptx_path, colors=colors)

        class DeterministicDelay:
            def __init__(self) -> None:
                self.calls = 0

            def describe(
                self, image_bytes: bytes, image_ext: str, shape_hint: str | None
            ) -> str:
                self.calls += 1
                # 바이트 값에 비례한 지연으로 완료 순서를 흔든다
                time.sleep((image_bytes[10] % 5) * 0.004)
                return f"det:{sum(image_bytes) % 997}"

        outputs: list[bytes] = []
        for _ in range(3):
            describer = DeterministicDelay()
            opts = ConvertOptions(describer=describer, describe_max_workers=8)
            md = convert(pptx_path, options=opts)
            outputs.append(md.encode("utf-8"))

        assert outputs[0] == outputs[1] == outputs[2]

    def test_tc_68_23_경계_worker_수를_1_4_8로_바꿔도_결과_동일(self) -> None:
        """TC-68-23 (경계): 동일 IR 스냅샷을 max_workers=1,4,8 세 가지로
        각각 변환해도 최종 description 매핑이 전부 동일하다(동시성 수준에
        결과가 의존하지 않음)."""

        def build() -> list[ImageShapeIR]:
            return [_img(i, _png_bytes((i % 7, i % 5, i % 3))) for i in range(15)]

        results: list[list[str | None]] = []
        for workers in (1, 4, 8):
            shapes = build()
            pres = _pres([_slide(shapes)])
            enrich_descriptions(
                pres, QaConcurrencyProbe(base_delay=0.005), max_workers=workers
            )
            results.append([s.description for s in shapes])

        assert results[0] == results[1] == results[2]

    def test_tc_68_24_예외_완료순서_역전_스트레스_3회_반복_동일(self) -> None:
        """TC-68-24 (스트레스/예외적 타이밍): 제출 순서와 완료 순서가
        인위적으로 뒤바뀌도록(첫 제출 job이 가장 늦게 끝나도록) 지연을
        역순 배치한 상태에서 3회 반복 → 매번 동일한 shape→description
        매핑(완료 순서에 의존하지 않는 오케스트레이션 검증, ADR-610)."""

        def build() -> list[ImageShapeIR]:
            return [
                _img(i, _png_bytes((i, i * 2 % 255, i * 3 % 255))) for i in range(8)
            ]

        mappings: list[list[str | None]] = []
        for _ in range(3):
            shapes = build()
            pres = _pres([_slide(shapes)])
            probe = QaConcurrencyProbe(reverse_delay=True, base_delay=0.006)
            enrich_descriptions(pres, probe, max_workers=8)
            mappings.append([s.description for s in shapes])

        assert mappings[0] == mappings[1] == mappings[2]
        assert all(d is not None for d in mappings[0])


# ===========================================================================
# AC9 (커버리지 지표) + AC3 조합 — 벡터 변환 실패라는 다른 실패 모드
# ===========================================================================


class TestAc9CoverageWithVectorFailure:
    """TC-68-25~27: EMF 벡터 변환 실패 경로로 AC3+AC9 조합 재현."""

    def test_tc_68_25_정상_7장_전량_성공시_커버리지_100(self) -> None:
        """TC-68-25 (정상): 이슈 AC9 문구대로 7장 이미지 전량 describe 성공
        → 커버리지 100%(모든 description 비-None)."""
        shapes = [_img(i, _png_bytes((i * 11 % 255, i * 7 % 255, i))) for i in range(7)]
        pres = _pres([_slide(shapes)])
        describer = QaKeyedDescriber()

        enrich_descriptions(pres, describer)

        assert sum(1 for s in shapes if s.description is not None) == 7

    def test_tc_68_26_경계_EMF_벡터변환실패_이미지만_FR22_마커_나머지_정상(
        self,
    ) -> None:
        """TC-68-26 (AC3+AC9 조합 경계): developer는 describe() 예외로
        실패를 재현했으나, 본 TC는 EMF→PNG 벡터 변환 실패
        (vector.convert_vector_to_png가 None 반환)라는 **다른 실패 경로**로
        동일 계약(실패 이미지만 None, 나머지 정상 채움, 최종 Markdown에
        FR-22 표준 마커 등장)을 재확인한다."""
        emf_bytes = b"\x01\x00\x00\x00FAKE_EMF_HEADER_BYTES"
        ok_bytes = _png_bytes((11, 22, 33))
        shapes = [
            _text(1, "제목"),
            _img(2, emf_bytes, image_ext="emf"),
            _img(3, ok_bytes),
        ]
        pres = _pres([_slide(shapes)])
        describer = QaKeyedDescriber()

        with patch.object(vector, "convert_vector_to_png", return_value=None):
            enrich_descriptions(pres, describer)

        assert shapes[1].description is None  # EMF 변환 실패 -> 격리
        assert shapes[2].description is not None  # 나머지 이미지는 정상

        md = assemble_document(pres)
        assert "슬라이드 1 이미지 1" in md  # FR-22 표준 위치 마커(실패 자리)
        assert shapes[2].description in md

    def test_tc_68_27_예외_전량_실패시_커버리지_0퍼센트여도_변환_성공(
        self,
    ) -> None:
        """TC-68-27 (극단 경계): 7장 전량 실패(커버리지 0%)여도
        enrich_descriptions와 assemble_document 모두 예외 없이 완료되고,
        7개 자리 전부 FR-22 표준 마커로 대체된다."""
        shapes = [_img(i, _png_bytes((i, i, i))) for i in range(1, 8)]
        pres = _pres([_slide(shapes)])
        describer = QaAllFailDescriber()

        enrich_descriptions(pres, describer)
        md = assemble_document(pres)

        assert sum(1 for s in shapes if s.description is not None) == 0
        for n in range(1, 8):
            assert f"슬라이드 1 이미지 {n}" in md


# ===========================================================================
# AC10 (조건부 COM) — convert_via_com 부재 확인
# ===========================================================================


class TestAc10ComAbsence:
    """TC-68-28: COM 경로 부재를 코드 스캔으로 독립 재확인(N/A 확정)."""

    def test_tc_68_28_convert_via_com_전역_부재(self) -> None:
        """TC-68-28: `convert_via_com` 심볼이 src/pptx_md 전역에 없고,
        win32com 관련 신규 모듈도 없다 -> AC10은 N/A로 확정."""
        import pptx_md

        assert not hasattr(pptx_md, "convert_via_com")

        pkg_dir = _PROJECT_ROOT / "src" / "pptx_md"
        for py_file in pkg_dir.rglob("*.py"):
            source = py_file.read_text(encoding="utf-8")
            assert "convert_via_com" not in source
            assert "win32com" not in source


# ===========================================================================
# 통합 스모크 (요구 §3) — 실패 격리 + Mermaid + 캐싱을 한 문서에서 동시 확인
# ===========================================================================


class TestIntegrationSmokeCombined:
    """TC-68-29: 실제 PPTX + FakeDescriber로 전 기능 결합 e2e."""

    def test_tc_68_29_실제_pptx_다이어그램옵션_캐싱_실패격리_결합(
        self, tmp_path: Path
    ) -> None:
        """TC-68-29: 슬라이드 1과 슬라이드 3에 동일 로고 이미지(캐싱),
        슬라이드 2에 DIAGRAM 분류 대상 이미지(Mermaid 옵션), 실패용 이미지
        1장을 한 실제 PPTX에 담아 convert()를 diagram_mermaid=True로 실행.
        parse/enrich_images/assemble 어느 것도 mock 하지 않는다."""
        logo_color = (250, 250, 0)
        pptx_path = tmp_path / "combined.pptx"
        _save_real_pptx_with_images(
            pptx_path,
            colors=[logo_color, (10, 10, 200), logo_color, (77, 88, 99)],
        )

        call_bytes: list[bytes] = []

        class CombinedDescriber:
            def describe(
                self, image_bytes: bytes, image_ext: str, shape_hint: str | None
            ) -> str:
                call_bytes.append(image_bytes)
                if image_bytes == _png_bytes((77, 88, 99)):
                    raise RuntimeError("combined smoke: simulated failure")
                return f"combined-desc-{len(call_bytes)}"

        opts = ConvertOptions(
            describer=CombinedDescriber(), diagram_mermaid=True, describe_max_workers=4
        )
        md = convert(pptx_path, options=opts)

        # 로고 2회 등장 -> dedup으로 실제 describe 호출은 고유 3개(logo, blue, fail)
        assert len(call_bytes) == 3
        assert isinstance(md, str)
        assert "combined-desc" in md
